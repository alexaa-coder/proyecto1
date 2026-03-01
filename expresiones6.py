# convert_multi_format_v5.py
# Conversor modular multi-formato para Docusaurus
# Soporta: PDF, DOCX, XLSX, XLS, XLSM, PPTX, TXT, PNG, ZIP
# 100% GRATIS - Sin APIs de pago
#
# NOVEDADES v5:
#   - Detección automática de PDF escaneado con OCR
#   - Procesamiento recursivo de archivos ZIP
#   - Manejo robusto de errores para migración masiva (600+ docs)
#   - Log de errores en archivo errors.log
#   - Normalización de nombres de archivo
#   - Soporte para extensiones en mayúsculas/minúsculas
#   - Métricas detalladas (éxitos, fallos, OCR usados)

import re
import time
import unicodedata
import logging
import tempfile
import shutil
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Optional, Tuple, List
from datetime import datetime
import warnings

# Dependencias según disponibilidad
try:
    import fitz  # PyMuPDF
    import pymupdf4llm
    PYMUPDF_DISPONIBLE = True
except ImportError:
    PYMUPDF_DISPONIBLE = False
    warnings.warn("pymupdf4llm no disponible - Los PDFs no se podrán procesar")

try:
    from docx import Document
    DOCX_DISPONIBLE = True
except ImportError:
    DOCX_DISPONIBLE = False
    warnings.warn("python-docx no disponible - Los DOCX no se podrán procesar")

try:
    import pandas as pd
    PANDAS_DISPONIBLE = True
except ImportError:
    PANDAS_DISPONIBLE = False
    warnings.warn("pandas no disponible - Los archivos Excel no se podrán procesar")

try:
    from pptx import Presentation
    PPTX_DISPONIBLE = True
except ImportError:
    PPTX_DISPONIBLE = False
    warnings.warn("python-pptx no disponible - Los PPTX no se podrán procesar")

try:
    from PIL import Image
    import pytesseract
    OCR_DISPONIBLE = True
except ImportError:
    OCR_DISPONIBLE = False
    warnings.warn("PIL/pytesseract no disponibles - OCR no estará disponible")


# ===========================================================================
# CONFIGURACIÓN GLOBAL
# ===========================================================================

_DEBUG_MODE = False
_OCR_MODE = False
_stats = None
_logger = None


def set_debug_mode(enabled: bool):
    """Activa o desactiva el modo debug globalmente."""
    global _DEBUG_MODE
    _DEBUG_MODE = enabled


def set_ocr_mode(enabled: bool):
    """Activa o desactiva el modo OCR globalmente."""
    global _OCR_MODE
    _OCR_MODE = enabled


def debug_log(mensaje: str, **kwargs):
    """Imprime mensajes solo si el modo debug está activo."""
    if _DEBUG_MODE:
        if kwargs:
            texto = mensaje.format(**kwargs)
        else:
            texto = mensaje
        print(f"[DEBUG] {texto}")


def setup_logger(output_dir: Path):
    """Configura logger para registrar errores en archivo."""
    global _logger
    log_file = output_dir / "errors.log"
    
    _logger = logging.getLogger('converter')
    _logger.setLevel(logging.ERROR)
    
    # Handler para archivo
    fh = logging.FileHandler(log_file, encoding='utf-8')
    fh.setLevel(logging.ERROR)
    
    # Formato
    formatter = logging.Formatter(
        '%(asctime)s - %(levelname)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    fh.setFormatter(formatter)
    
    _logger.addHandler(fh)
    
    return log_file


# ===========================================================================
# ESTADÍSTICAS DE CONVERSIÓN
# ===========================================================================

class EstadisticasConversion:
    """Almacena métricas de conversión para modo debug."""
    
    def __init__(self):
        self.reset()
    
    def reset(self):
        self.lineas_raw = 0
        self.lineas_limpias = 0
        self.chars_raw = 0
        self.chars_limpios = 0
        self.tablas_detectadas = 0
        self.tablas_convertidas = 0
        self.bloques_duplicados_eliminados = 0
        self.palabras_guionadas_unidas = 0
        self.referencias_normalizadas = 0
        self.archivos_por_tipo = {}
        self.archivos_con_ocr = 0
        self.pdfs_escaneados = 0
        self.zips_procesados = 0
        self.archivos_en_zips = 0
        self.tiempo_total = 0.0
    
    def incrementar_tipo(self, extension: str):
        """Cuenta archivos procesados por tipo."""
        self.archivos_por_tipo[extension] = self.archivos_por_tipo.get(extension, 0) + 1
    
    def mostrar_resumen(self):
        """Imprime resumen de estadísticas en modo debug."""
        if not _DEBUG_MODE:
            return
        
        print("\n" + "─" * 60)
        print("📊 ESTADÍSTICAS DE CONVERSIÓN")
        print("─" * 60)
        
        if self.archivos_por_tipo:
            print("📁 Archivos por tipo:")
            for ext, count in sorted(self.archivos_por_tipo.items()):
                print(f"   {ext:10s} : {count:3d} archivos")
        
        if self.pdfs_escaneados > 0:
            print(f"\n📄 PDFs escaneados detectados: {self.pdfs_escaneados}")
        
        if self.archivos_con_ocr > 0:
            print(f"🔍 Archivos procesados con OCR: {self.archivos_con_ocr}")
        
        if self.zips_procesados > 0:
            print(f"\n📦 Archivos ZIP:")
            print(f"   ZIPs procesados     : {self.zips_procesados}")
            print(f"   Archivos extraídos  : {self.archivos_en_zips}")
        
        if self.lineas_raw > 0:
            print(f"\n📄 Líneas:")
            print(f"   Raw    : {self.lineas_raw:,}")
            print(f"   Limpias: {self.lineas_limpias:,}")
            print(f"   Reducción: {self.lineas_raw - self.lineas_limpias:,} líneas")
        
        if self.chars_raw > 0:
            print(f"\n📝 Caracteres:")
            print(f"   Raw    : {self.chars_raw:,}")
            print(f"   Limpios: {self.chars_limpios:,}")
            reduccion_pct = (self.chars_raw - self.chars_limpios) / self.chars_raw * 100 if self.chars_raw else 0
            print(f"   Reducción: {self.chars_raw - self.chars_limpios:,} chars ({reduccion_pct:.1f}%)")

        if self.tablas_detectadas > 0:
            print(f"\n📊 Tablas:")
            print(f"   Detectadas  : {self.tablas_detectadas}")
            print(f"   Convertidas : {self.tablas_convertidas}")

        if self.bloques_duplicados_eliminados > 0:
            print(f"\n📦 Bloques duplicados eliminados: {self.bloques_duplicados_eliminados}")

        if self.tiempo_total > 0:
             print(f"\n⏱  Tiempo total de ejecución: {self.tiempo_total:.2f}s")
        
        print("─" * 60 + "\n")


# ===========================================================================
# UTILIDADES
# ===========================================================================

def normalizar_nombre_archivo(nombre: str) -> str:
    """
    Normaliza nombre de archivo para evitar problemas de caracteres especiales.
    Convierte a slug seguro manteniendo extensión.
    """
    # Separar nombre y extensión
    ruta = Path(nombre)
    stem = ruta.stem
    ext = ruta.suffix.lower()
    
    # Normalizar el nombre (sin extensión)
    stem_normalizado = unicodedata.normalize("NFKD", stem)
    stem_normalizado = stem_normalizado.encode("ascii", "ignore").decode("ascii")
    stem_normalizado = re.sub(r"[^\w\-]", "-", stem_normalizado.lower()).strip("-")
    
    # Evitar nombres vacíos
    if not stem_normalizado:
        stem_normalizado = "documento"
    
    return stem_normalizado + ext


def simple_slug(text: str) -> str:
    """Genera slug seguro para IDs de Docusaurus."""
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("ascii", "ignore").decode("ascii")
    return re.sub(r"[^\w\-]", "-", text.lower()).strip("-")


def get_frontmatter(file_path: Path) -> str:
    """Genera frontmatter YAML para Docusaurus."""
    doc_id = simple_slug(file_path.stem)
    title = file_path.stem.replace("-", " ").replace("_", " ").title()
    return f"""---
id: {doc_id}
title: "{title}"
sidebar_label: "{title}"
---

"""


# ===========================================================================
# PIPELINE DE LIMPIEZA MARKDOWN
# ===========================================================================

def unir_palabras_guionadas(texto: str) -> str:
    """Restaura palabras cortadas con guión de fin de línea."""
    contador = len(re.findall(r"([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])-\n([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])", texto))
    _stats.palabras_guionadas_unidas += contador
    debug_log("Uniendo {n} palabras guionadas", n=contador)
    return re.sub(r"([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])-\n([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])", r"\1\2", texto)


def limpiar_indice(texto: str) -> str:
    """Limpia entradas de índices."""
    RE_TIENE_DOTS = re.compile(r"\.{3,}\s*\d")
    RE_QUITAR_DOTS_NUM = re.compile(r"\s*\.{3,}\s*\d+")
    
    lineas = texto.split("\n")
    resultado = []
    lineas_eliminadas = 0
    
    for linea in lineas:
        if not RE_TIENE_DOTS.search(linea):
            resultado.append(linea)
            continue
        
        limpia = RE_QUITAR_DOTS_NUM.sub("", linea).strip()
        if limpia:
            resultado.append(limpia)
        else:
            lineas_eliminadas += 1
    
    debug_log("Índice: {n} líneas eliminadas", n=lineas_eliminadas)
    return "\n".join(resultado)


def limpiar_encabezados_pies(texto: str) -> str:
    """Elimina encabezados y pies de página."""
    patron_num_norma = r"^\s*\d+\s*[-–—]\s*[A-Za-z]{2,}[\w/\-\s.:]+\d{4,}\s*$"
    patron_norma_num = r"^[A-Za-z]{2,}[\w/\-\s.:]+\d{4,}\s*[-–—]\s*\d+\s*$"
    patron_pagina = r"^\s*\d+\s*$"
    
    lineas = texto.split("\n")
    lineas_eliminadas = 0
    resultado = []
    
    for linea in lineas:
        if (re.match(patron_num_norma, linea) or
            re.match(patron_norma_num, linea) or
            re.match(patron_pagina, linea)):
            lineas_eliminadas += 1
            continue
        resultado.append(linea)
    
    debug_log("Encabezados/pies: {n} eliminados", n=lineas_eliminadas)
    return "\n".join(resultado)


def fusionar_parrafos_rotos(texto: str) -> str:
    """Fusiona líneas cortadas arbitrariamente."""
    FIN_FRASE = (".", ":", ";", "!", "?", ")")
    MARCADORES = ("#", "-", "*", "|", ">", "```")
    
    lineas = texto.split("\n")
    lineas_fusionadas = 0
    resultado = []
    i = 0
    
    while i < len(lineas):
        linea = lineas[i].rstrip()
        
        if not linea:
            resultado.append("")
            i += 1
            continue
        
        if linea.startswith(MARCADORES) or re.match(r"^\d+(?:\.\d+)*\s", linea):
            resultado.append(linea)
            i += 1
            continue
        
        if linea.endswith(FIN_FRASE):
            resultado.append(linea)
            i += 1
            continue
        
        if i + 1 < len(lineas):
            sig = lineas[i + 1].lstrip()
            if not sig or sig.startswith(MARCADORES):
                resultado.append(linea)
                i += 1
            else:
                resultado.append(linea + " " + sig)
                lineas_fusionadas += 1
                i += 2
        else:
            resultado.append(linea)
            i += 1
    
    debug_log("Párrafos: {n} fusionados", n=lineas_fusionadas)
    return "\n".join(resultado)


def normalizar_referencias(texto: str) -> str:
    """Limpia referencias cruzadas."""
    contador = 0
    
    nuevo, n1 = re.subn(
        r"(\b(?:véase|see|clause|cláusula)\b)\n(\d[\d\.]*)",
        r"\1 \2",
        texto,
        flags=re.IGNORECASE,
    )
    contador += n1
    
    nuevo, n2 = re.subn(
        r"\[([A-Z]{2,}[/\-\w\s]+?)\s*:\s*(\d{4})\]",
        r"[\1:\2]",
        nuevo,
    )
    contador += n2
    
    _stats.referencias_normalizadas += contador
    debug_log("Referencias: {n} normalizadas", n=contador)
    return nuevo


def arreglar_html_mdx(texto: str) -> str:
    """Convierte etiquetas HTML a JSX y escapa caracteres peligrosos fuera de tags."""
    # 1. Arreglar tags conocidos para que sean MDX-safe
    texto = re.sub(r"<br\s*>", "<br />", texto, flags=re.IGNORECASE)
    texto = re.sub(r"<hr\s*>", "<hr />", texto, flags=re.IGNORECASE)
    texto = re.sub(r"<img([^>]+)(?<!/)>", r"<img\1 />", texto, flags=re.IGNORECASE)
    
    # 2. Escapar símbolos que MDX confunde con JSX (solo si no parecen tags cerrados)
    # Escapar < si no es el inicio de un tag conocido o cierre de tag
    # Esto evita errores como "Unexpected character 1 before name" en "< 100"
    texto = re.sub(r"<(?!/?(?:p|br|hr|img|div|span|strong|em|b|i|u|h\d)\b)", "&lt;", texto, flags=re.IGNORECASE)
    
    # 3. Escapar llaves literales (interfieren con expresiones JS en MDX)
    texto = texto.replace("{", "&#123;").replace("}", "&#125;")
    
    return texto


def limpiar_espacios_multiples(texto: str) -> str:
    """Normaliza espacios."""
    texto = re.sub(r" {2,}", " ", texto)
    texto = re.sub(r"\t", "    ", texto)
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    return texto


def justificar_texto(texto: str) -> str:
    """
    Envuelve los párrafos extraídos en etiquetas de justificación para Docusaurus/MDX.
    Usa expresiones regulares para identificar y transformar los bloques.
    """
    debug_log("Aplicando justificación de texto (modo regex)...")
    
    # Patrón para identificar bloques estructurales que NO deben justificarse
    # (Headers, Listas, Tablas, Bloques de código, Frontmatter, tags HTML, etc.)
    PATRON_ESTRUCTURAL = r"^(?:#|\-|\*|\d+\.|\||```|<img|<p|<div|<br|<hr|&lt;|&#123;|--|id:|title:|sidebar_label:)"
    
    def callback_justificar(match):
        bloque = match.group(0)
        content = bloque.strip()
        
        if not content:
            return bloque
            
        # Si el bloque contiene alguna línea estructural o es muy corto, no justificar
        lines = content.split('\n')
        if any(re.match(PATRON_ESTRUCTURAL, l.strip(), re.IGNORECASE) for l in lines) or len(content) < 50:
            return bloque
            
        # Justificar bloque usando sintaxis compatible con Docusaurus/MDX
        return f'<div style={{{{textAlign: "justify"}}}}>\n\n{content}\n\n</div>'

    # Regex para capturar bloques: cualquier secuencia que no contenga una doble línea vacía
    regex_bloques = r"((?:(?!\n\n).)+)"
    
    # Aplicar transformación a cada bloque capturado
    return re.sub(regex_bloques, callback_justificar, texto, flags=re.DOTALL)


def limpiar_markdown_completo(md_texto: str, verbose: bool = True) -> str:
    """Pipeline completo de limpieza Markdown."""
    _stats.chars_raw += len(md_texto)
    _stats.lineas_raw += md_texto.count("\n") + 1
    
    debug_log("="*60)
    debug_log("INICIANDO PIPELINE DE LIMPIEZA")
    debug_log("="*60)
    
    pasos = [
        ("🔗 Uniendo palabras guionadas", unir_palabras_guionadas),
        ("🧹 Limpiando índice", limpiar_indice),
        ("🧹 Eliminando encabezados/pies", limpiar_encabezados_pies),
        ("🔗 Normalizando referencias", normalizar_referencias),
        ("🧹 Fusionando párrafos", fusionar_parrafos_rotos),
        ("🧹 Limpiando espacios", limpiar_espacios_multiples),
        ("🧹 Arreglando HTML/MDX", arreglar_html_mdx),
        ("⚖️  Justificando texto", justificar_texto),
    ]
    
    texto = md_texto
    for nombre, fn in pasos:
        if verbose:
            print(f"      {nombre}...")
        texto = fn(texto)
    
    _stats.chars_limpios += len(texto)
    _stats.lineas_limpias += texto.count("\n") + 1
    
    return texto.strip()


# ===========================================================================
# HANDLERS POR TIPO DE ARCHIVO
# ===========================================================================

def aplicar_ocr_a_pdf(pdf_path: Path) -> Tuple[Optional[str], bool]:
    """
    Aplica OCR a un PDF completo usando Tesseract.
    """
    if not OCR_DISPONIBLE:
        debug_log("OCR no disponible - pytesseract no instalado")
        return None, False
    
    try:
        debug_log("Aplicando OCR con Tesseract...")
        doc = fitz.open(pdf_path)
        textos_paginas = []
        
        for page_num, page in enumerate(doc, 1):
            debug_log(f"OCR en página {page_num}/{len(doc)}...")
            
            # Renderizar página como imagen
            pix = page.get_pixmap(dpi=300)
            img_data = pix.tobytes("png")
            
            # Convertir a PIL Image
            from io import BytesIO
            img = Image.open(BytesIO(img_data))
            
            # Extraer texto con Tesseract
            texto = pytesseract.image_to_string(img, lang='spa')
            textos_paginas.append(texto)
        
        doc.close()
        texto_completo = "\n\n".join(textos_paginas)
        debug_log(f"OCR completado: {len(texto_completo)} caracteres")
        
        _stats.archivos_con_ocr += 1
        
        return texto_completo, True
    
    except Exception as e:
        debug_log(f"Error en OCR: {e}")
        return None, False


def handle_pdf(file_path: Path) -> Optional[str]:
    """
    Convierte PDF a Markdown.
    """
    if not PYMUPDF_DISPONIBLE:
        print(f"   ⚠️  pymupdf4llm no disponible")
        return None
    
    try:
        debug_log("Analizando PDF...")
        
        # Paso 1: Intentar extracción normal
        doc = fitz.open(file_path)
        num_paginas = len(doc)
        
        # Extraer texto de todas las páginas para análisis
        texto_completo = ""
        for page in doc:
            texto_completo += page.get_text()
        
        doc.close()
        
        chars_promedio_por_pagina = len(texto_completo) / num_paginas if num_paginas > 0 else 0
        
        debug_log(f"PDF: {num_paginas} páginas, {chars_promedio_por_pagina:.0f} chars/pág promedio")
        
        # Paso 2: Detectar si es PDF escaneado
        es_escaneado = chars_promedio_por_pagina < 500
        
        if es_escaneado:
            debug_log("⚠️  PDF posiblemente escaneado (poco texto extraíble)")
            _stats.pdfs_escaneados += 1
            
            # Si OCR está habilitado, aplicarlo
            if _OCR_MODE:
                print(f"   🔍 Aplicando OCR (PDF escaneado)...")
                texto_ocr, ocr_ok = aplicar_ocr_a_pdf(file_path)
                
                if ocr_ok and texto_ocr:
                    debug_log(f"OCR exitoso: {len(texto_ocr)} caracteres")
                    return texto_ocr
                else:
                    print(f"   ⚠️  OCR falló, usando extracción normal")
        
        # Paso 3: Extracción normal con pymupdf4llm
        debug_log("Extrayendo con pymupdf4llm...")
        md_raw = pymupdf4llm.to_markdown(str(file_path))
        debug_log(f"Extraído: {len(md_raw)} caracteres")
        
        return md_raw
    
    except Exception as e:
        print(f"   ❌ Error: {e}")
        if _logger:
            _logger.error(f"PDF {file_path.name}: {e}")
        return None


def handle_docx(file_path: Path) -> Optional[str]:
    """Convierte DOCX a Markdown."""
    if not DOCX_DISPONIBLE:
        print(f"   ⚠️  python-docx no disponible")
        return None
    
    try:
        debug_log("Extrayendo DOCX...")
        doc = Document(file_path)
        
        markdown_lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                markdown_lines.append("")
                continue
            
            if para.style.name.startswith('Heading'):
                level = para.style.name.replace('Heading ', '')
                try:
                    level_num = int(level)
                    markdown_lines.append(f"{'#' * level_num} {text}")
                except:
                    markdown_lines.append(text)
            else:
                markdown_lines.append(text)
        
        for table in doc.tables:
            _stats.tablas_detectadas += 1
            markdown_lines.append("")
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                markdown_lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    markdown_lines.append("|" + "---|" * len(cells))
            markdown_lines.append("")
            _stats.tablas_convertidas += 1
        
        md_raw = "\n".join(markdown_lines)
        debug_log(f"DOCX: {len(md_raw)} caracteres")
        return md_raw
    
    except Exception as e:
        print(f"   ❌ Error: {e}")
        if _logger:
            _logger.error(f"DOCX {file_path.name}: {e}")
        return None


def handle_excel(file_path: Path) -> Optional[str]:
    """Convierte Excel a Markdown."""
    if not PANDAS_DISPONIBLE:
        print(f"   ⚠️  pandas no disponible")
        return None
    
    try:
        debug_log("Extrayendo Excel...")
        excel_file = pd.ExcelFile(file_path)
        markdown_lines = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            markdown_lines.append(f"\n## {sheet_name}\n")
            
            _stats.tablas_detectadas += 1
            if not df.empty:
                headers = df.columns.tolist()
                markdown_lines.append("| " + " | ".join(str(h) for h in headers) + " |")
                markdown_lines.append("|" + "---|" * len(headers))
                
                for _, row in df.iterrows():
                    cells = [str(val) if pd.notna(val) else "" for val in row]
                    markdown_lines.append("| " + " | ".join(cells) + " |")
                
                markdown_lines.append("")
                _stats.tablas_convertidas += 1
        
        md_raw = "\n".join(markdown_lines)
        debug_log(f"Excel: {len(excel_file.sheet_names)} hojas, {len(md_raw)} chars")
        return md_raw
    
    except Exception as e:
        print(f"   ❌ Error: {e}")
        if _logger:
            _logger.error(f"Excel {file_path.name}: {e}")
        return None


def handle_pptx(file_path: Path) -> Optional[str]:
    """Convierte PowerPoint a Markdown."""
    if not PPTX_DISPONIBLE:
        print(f"   ⚠️  python-pptx no disponible")
        return None
    
    try:
        debug_log("Extrayendo PPTX...")
        prs = Presentation(file_path)
        markdown_lines = []
        
        for i, slide in enumerate(prs.slides, 1):
            markdown_lines.append(f"\n## Diapositiva {i}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    markdown_lines.append(shape.text.strip())
                    markdown_lines.append("")
                
                if shape.has_table:
                    _stats.tablas_detectadas += 1
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        markdown_lines.append("| " + " | ".join(cells) + " |")
                        if row_idx == 0:
                            markdown_lines.append("|" + "---|" * len(cells))
                    markdown_lines.append("")
                    _stats.tablas_convertidas += 1
        
        md_raw = "\n".join(markdown_lines)
        debug_log(f"PPTX: {len(prs.slides)} diapositivas, {len(md_raw)} chars")
        return md_raw
    
    except Exception as e:
        print(f"   ❌ Error: {e}")
        if _logger:
            _logger.error(f"PPTX {file_path.name}: {e}")
        return None


def handle_txt(file_path: Path) -> Optional[str]:
    """Lee archivo TXT."""
    try:
        debug_log("Leyendo TXT...")
        md_raw = file_path.read_text(encoding='utf-8')
        debug_log(f"TXT: {len(md_raw)} caracteres")
        return md_raw
    except UnicodeDecodeError:
        try:
            md_raw = file_path.read_text(encoding='latin-1')
            debug_log(f"TXT (latin-1): {len(md_raw)} caracteres")
            return md_raw
        except Exception as e:
            print(f"   ❌ Error: {e}")
            if _logger:
                _logger.error(f"TXT {file_path.name}: {e}")
            return None
    except Exception as e:
        print(f"   ❌ Error: {e}")
        if _logger:
            _logger.error(f"TXT {file_path.name}: {e}")
        return None


def handle_png(file_path: Path) -> Optional[str]:
    """
    Extrae texto de imagen PNG usando OCR.
    """
    if not OCR_DISPONIBLE:
        print(f"   ⚠️  pytesseract no disponible")
        # Insertar como imagen
        md_raw = f"![{file_path.stem}](./{file_path.name})\n\n*Imagen: {file_path.name}*"
        debug_log(f"PNG insertado como imagen (OCR no disponible)")
        return md_raw
    
    try:
        debug_log("Aplicando OCR a PNG...")
        image = Image.open(file_path)
        texto = pytesseract.image_to_string(image, lang='spa')
        
        # Verificar si se extrajo texto significativo
        texto_limpio = texto.strip()
        chars_detectados = len(texto_limpio)
        
        debug_log(f"PNG OCR: {chars_detectados} caracteres detectados")
        
        if chars_detectados > 50:  # Umbral mínimo de texto
            debug_log("PNG procesado como texto")
            _stats.archivos_con_ocr += 1
            return texto_limpio
        else:
            # Poco o ningún texto → insertar como imagen
            debug_log("PNG insertado como imagen (poco texto detectado)")
            md_raw = f"![{file_path.stem}](./{file_path.name})\n\n*Imagen: {file_path.name}*"
            return md_raw
    
    except Exception as e:
        print(f"   ❌ Error: {e}")
        if _logger:
            _logger.error(f"PNG {file_path.name}: {e}")
        # En caso de error, insertar como imagen
        md_raw = f"![{file_path.stem}](./{file_path.name})\n\n*Imagen: {file_path.name}*"
        return md_raw


def handle_zip(file_path: Path, output_base_dir: Path) -> Tuple[int, int]:
    """
    Descomprime ZIP y procesa recursivamente su contenido.
    """
    import zipfile
    
    try:
        debug_log(f"Procesando ZIP: {file_path.name}")
        
        # Crear carpeta temporal
        temp_dir = tempfile.mkdtemp(prefix="zip_extract_")
        temp_path = Path(temp_dir)
        
        debug_log(f"Extrayendo en: {temp_path}")
        
        # Extraer ZIP
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_path)
            archivos_extraidos = zip_ref.namelist()
        
        num_archivos = len(archivos_extraidos)
        debug_log(f"ZIP: {num_archivos} archivos extraídos")
        
        _stats.zips_procesados += 1
        _stats.archivos_en_zips += num_archivos
        
        # Procesar recursivamente
        exitos = 0
        fallos = 0
        
        # Buscar todos los archivos soportados en el ZIP
        for archivo_extraido in temp_path.rglob("*"):
            if archivo_extraido.is_file():
                ext = archivo_extraido.suffix.lower()
                
                if ext in HANDLERS:
                    # Determinar ruta de salida manteniendo estructura
                    rel_path = archivo_extraido.relative_to(temp_path)
                    output_path = output_base_dir / file_path.stem / rel_path.with_suffix('.md')
                    
                    debug_log(f"  Procesando: {archivo_extraido.name}")
                    
                    if convertir_archivo(archivo_extraido, output_path, verbose=False):
                        exitos += 1
                    else:
                        fallos += 1
        
        # Limpiar carpeta temporal
        shutil.rmtree(temp_path)
        debug_log(f"ZIP procesado: {exitos} éxitos, {fallos} fallos")
        
        return exitos, fallos
    
    except Exception as e:
        print(f"   ❌ Error procesando ZIP: {e}")
        if _logger:
            _logger.error(f"ZIP {file_path.name}: {e}")
        return 0, 1


# ===========================================================================
# DISPATCHER DE HANDLERS
# ===========================================================================

HANDLERS = {
    '.pdf': handle_pdf,
    '.docx': handle_docx,
    '.xlsx': handle_excel,
    '.xls': handle_excel,
    '.xlsm': handle_excel,
    '.pptx': handle_pptx,
    '.txt': handle_txt,
    '.png': handle_png,
}

EXTENSIONES_IGNORADAS = {'.exe', '.url'}


def convertir_archivo(file_path: Path, output_path: Path, verbose: bool = True) -> bool:
    """
    Convierte un archivo a Markdown según su extensión.
    """
    t_file_start = time.perf_counter()

    # Normalizar extensión (soportar mayúsculas)
    extension = file_path.suffix.lower()
    
    # Verificar extensiones ignoradas
    if extension in EXTENSIONES_IGNORADAS:
        if verbose:
            print(f"\n⚠️  Ignorando: {file_path.name} (extensión {extension})")
        return False
    
    # Caso especial: ZIP
    if extension == '.zip':
        if verbose:
            print(f"\n📦 Procesando ZIP: {file_path.name}")
        exitos, fallos = handle_zip(file_path, output_path.parent)
        return exitos > 0
    
    # Verificar handler
    if extension not in HANDLERS:
        if verbose:
            print(f"\n⚠️  Formato desconocido: {file_path.name} ({extension})")
        return False
    
    try:
        t_inicio = time.time()
        
        if verbose:
            print(f"\n📖 Procesando: {file_path.name}")
            print(f"   📦 Tipo: {extension.upper()[1:]}")
        
        # Obtener handler
        handler = HANDLERS[extension]
        
        # Extraer contenido
        if verbose:
            print(f"   ⚙️  Extrayendo contenido...")
        
        md_raw = handler(file_path)
        
        if md_raw is None:
            return False
        
        if verbose:
            print(f"   📊 Extraído: {len(md_raw):,} caracteres")
            print(f"   🔧 Limpiando Markdown...")
        
        # Limpiar
        md_limpio = limpiar_markdown_completo(md_raw, verbose=verbose)
        
        # Frontmatter
        contenido_final = get_frontmatter(file_path) + md_limpio
        
        # Normalizar nombre de salida
        output_path_normalizado = output_path.parent / normalizar_nombre_archivo(output_path.name)
        
        # Guardar
        output_path_normalizado.parent.mkdir(parents=True, exist_ok=True)
        output_path_normalizado.write_text(contenido_final, encoding='utf-8')
        
        # Estadísticas
        _stats.incrementar_tipo(extension)
        
        t_total = time.time() - t_inicio
        
        if verbose:
            reduccion = len(md_raw) - len(md_limpio)
            pct = reduccion / len(md_raw) * 100 if md_raw else 0
            print(f"   ✅ Guardado: {output_path_normalizado}")
            print(f"   📉 Reducción: {reduccion:,} chars ({pct:.1f}%)")
        
        t_file_end = time.perf_counter()
        duracion_archivo = t_file_end - t_file_start
        debug_log("⏱  Tiempo archivo: {t:.2f}s", t=duracion_archivo)
        
        return True
    
    except Exception as e:
        print(f"   ❌ ERROR: {e}")
        if _logger:
            _logger.error(f"{file_path.name}: {str(e)}")
        if _DEBUG_MODE:
            import traceback
            traceback.print_exc()
        return False


# ===========================================================================
# PROCESAMIENTO MASIVO
# ===========================================================================

def procesar_carpeta(
    input_dir: Path,
    output_dir: Path,
    verbose: bool = True,
    workers: int = 1,
    extensiones_filtro: Optional[list] = None,
):
    """Procesa todos los archivos de una carpeta."""
    t_process_start = time.perf_counter()
    separador = "=" * 70
    
    if verbose:
        print(f"\n{separador}")
        print(f"🔄 CONVERSIÓN MULTI-FORMATO — v5 (MIGRACIÓN MASIVA)")
        print(f"{separador}")
        print(f"📁 Origen : {input_dir}")
        print(f"📂 Destino: {output_dir}")
        print(f"⚡ Workers: {workers}")
        if _DEBUG_MODE:
            print(f"🐛 Modo DEBUG: ACTIVO")
        if _OCR_MODE:
            print(f"🔍 Modo OCR: ACTIVO")
    
    # Setup logger
    log_file = setup_logger(output_dir)
    if verbose:
        print(f"📋 Log errores: {log_file}")
    
    # Buscar archivos
    archivos = []
    extensiones_buscar = extensiones_filtro or list(HANDLERS.keys()) + ['.zip']
    
    for ext in extensiones_buscar:
        archivos.extend(input_dir.rglob(f"*{ext}"))
        # Soportar extensiones en mayúsculas
        archivos.extend(input_dir.rglob(f"*{ext.upper()}"))
    
    # Eliminar duplicados
    archivos = list(set(archivos))
    
    if not archivos:
        print(f"❌ No se encontraron archivos.")
        print(f"   Extensiones buscadas: {', '.join(extensiones_buscar)}")
        return
    
    print(f"📊 Total archivos encontrados: {len(archivos)}\n")
    print(f"⏱  Inicio: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
    
    t_inicio = time.time()
    
    exitos = 0
    errores = 0
    ignorados = 0
    
    def _tarea(file_path: Path):
        rel_path = file_path.relative_to(input_dir)
        out_path = output_dir / rel_path.with_suffix('.md')
        return convertir_archivo(file_path, out_path, verbose=(workers == 1 and verbose))
    
    if workers == 1:
        # Secuencial
        for file_path in archivos:
            if file_path.suffix.lower() in EXTENSIONES_IGNORADAS:
                ignorados += 1
                continue
            
            if _tarea(file_path):
                exitos += 1
            else:
                errores += 1
    else:
        # Paralelo
        with ThreadPoolExecutor(max_workers=workers) as executor:
            futuros = {}
            for file_path in archivos:
                if file_path.suffix.lower() in EXTENSIONES_IGNORADAS:
                    ignorados += 1
                    continue
                futuros[executor.submit(_tarea, file_path)] = file_path
            
            procesados = 0
            for futuro in as_completed(futuros):
                procesados += 1
                nombre = futuros[futuro].name
                try:
                    ok = futuro.result()
                    estado = "✅" if ok else "❌"
                    print(f"  [{procesados}/{len(archivos)-ignorados}] {estado} {nombre}")
                    if ok:
                        exitos += 1
                    else:
                        errores += 1
                except Exception as e:
                    print(f"  [{procesados}/{len(archivos)-ignorados}] ❌ {nombre} - {e}")
                    errores += 1
    
    t_total = time.time() - t_inicio
    velocidad = (exitos + errores) / (t_total / 60) if t_total > 0 else 0
    
    t_process_end = time.perf_counter()
    _stats.tiempo_total = t_process_end - t_process_start

    # Mostrar estadísticas
    _stats.mostrar_resumen()
    
    print(f"\n{separador}")
    print(f"📊 RESUMEN FINAL")
    print(f"{separador}")
    print(f"✅ Conversiones exitosas : {exitos}")
    print(f"❌ Errores              : {errores}")
    print(f"⏭️  Ignorados             : {ignorados}")
    if _stats.archivos_con_ocr > 0:
        print(f"🔍 Procesados con OCR    : {_stats.archivos_con_ocr}")
    if _stats.pdfs_escaneados > 0:
        print(f"📄 PDFs escaneados       : {_stats.pdfs_escaneados}")
    print(f"📁 Total procesados     : {exitos + errores}")
    print(f"⏱  Tiempo total         : {t_total:.1f}s ({t_total/60:.1f} min)")
    print(f"⚡ Velocidad            : {velocidad:.1f} archivos/min")
    print(f"📂 Archivos en          : {output_dir}")
    if errores > 0:
        print(f"📋 Ver errores en       : {log_file}")
    print(f"🕐 Fin: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print(f"{separador}\n")


# ===========================================================================
# CLI
# ===========================================================================

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Conversor multi-formato a Markdown para Docusaurus — v5",
        epilog=(
            "Soportes: PDF, DOCX, XLSX, PPTX, TXT, PNG, ZIP\n"
            "Optimizado para migración masiva (600+ documentos)\n"
            "\n"
            "Ejemplos:\n"
            "  python convert_multi_format_v5.py\n"
            "  python convert_multi_format_v5.py --ocr\n"
            "  python convert_multi_format_v5.py --workers 4 --debug\n"
            "  python convert_multi_format_v5.py --only-pdf --ocr\n"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    
    parser.add_argument("--input", default="pdf",
                        help="Carpeta origen (default: pdf/)")
    parser.add_argument("--output", default="docs",
                        help="Carpeta destino (default: docs/)")
    parser.add_argument("--workers", default=1, type=int,
                        help="Hilos paralelos (default: 1)")
    parser.add_argument("--ocr", action="store_true",
                        help="Activar OCR para PDFs escaneados y PNGs")
    parser.add_argument("--debug", action="store_true",
                        help="Modo debug con estadísticas")
    parser.add_argument("--quiet", action="store_true",
                        help="Modo silencioso")
    parser.add_argument("--only-pdf", action="store_true")
    parser.add_argument("--only-docx", action="store_true")
    parser.add_argument("--only-excel", action="store_true")
    parser.add_argument("--only-pptx", action="store_true")
    
    args = parser.parse_args()
    
    # Activar modos
    if args.debug:
        set_debug_mode(True)
    if args.ocr:
        set_ocr_mode(True)
        if not OCR_DISPONIBLE:
            print("⚠️  OCR solicitado pero pytesseract no está disponible")
            print("   Instala: pip install pytesseract pillow")
            print("   Y Tesseract en el sistema")
    
    # Inicializar estadísticas
    _stats = EstadisticasConversion()
    
    # Filtros
    extensiones_filtro = None
    if args.only_pdf:
        extensiones_filtro = ['.pdf']
    elif args.only_docx:
        extensiones_filtro = ['.docx']
    elif args.only_excel:
        extensiones_filtro = ['.xlsx', '.xls', '.xlsm']
    elif args.only_pptx:
        extensiones_filtro = ['.pptx']
    
    # Validar carpetas
    input_path = Path(args.input)
    output_path = Path(args.output)
    
    if not input_path.exists():
        print(f"⚠️  Carpeta '{input_path}' no existe. Creándola...")
        input_path.mkdir(parents=True)
        print(f"   ℹ️  Coloca archivos en '{input_path}' y ejecuta de nuevo.")
        raise SystemExit(0)
    
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Procesar
    procesar_carpeta(
        input_path,
        output_path,
        verbose=not args.quiet,
        workers=args.workers,
        extensiones_filtro=extensiones_filtro,
    )
