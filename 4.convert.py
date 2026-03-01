# convert_multi_format_v4.py
# Conversor modular multi-formato para Docusaurus
# Soporta: PDF, DOCX, XLSX, XLS, XLSM, PPTX, TXT, PNG, ZIP
# 100% GRATIS - Sin APIs de pago
#
# ARQUITECTURA v4:
#   - Sistema modular de handlers por tipo de archivo
#   - Detección automática de formato
#   - Pipeline de limpieza Markdown reutilizable
#   - Soporte para procesamiento paralelo
#   - Modo debug integrado

import re
import time
import unicodedata
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Optional, Callable
import warnings

# Dependencias según disponibilidad
try:
    import pymupdf4llm
    PYMUPDF_DISPONIBLE = True
except ImportError:
    PYMUPDF_DISPONIBLE = False
    warnings.warn("pymupdf4llm no disponible - Los PDFs no se podrán procesar")

try:
    from docx import Document
    DOCX_DISPONIBLE = True
except ImportError:
    DOCX_DISPONIBLE = False
    warnings.warn("python-docx no disponible - Los DOCX no se podrán procesar")

try:
    import pandas as pd
    PANDAS_DISPONIBLE = True
except ImportError:
    PANDAS_DISPONIBLE = False
    warnings.warn("pandas no disponible - Los archivos Excel no se podrán procesar")

try:
    from pptx import Presentation
    PPTX_DISPONIBLE = True
except ImportError:
    PPTX_DISPONIBLE = False
    warnings.warn("python-pptx no disponible - Los PPTX no se podrán procesar")

try:
    from PIL import Image
    import pytesseract
    OCR_DISPONIBLE = True
except ImportError:
    OCR_DISPONIBLE = False
    warnings.warn("PIL/pytesseract no disponibles - Las imágenes no se podrán procesar con OCR")


# ===========================================================================
# CONFIGURACIÓN GLOBAL
# ===========================================================================

_DEBUG_MODE = False
_stats = None  # Se inicializa en main()


def set_debug_mode(enabled: bool):
    """Activa o desactiva el modo debug globalmente."""
    global _DEBUG_MODE
    _DEBUG_MODE = enabled


def debug_log(mensaje: str, **kwargs):
    """Imprime mensajes solo si el modo debug está activo."""
    if _DEBUG_MODE:
        if kwargs:
            texto = mensaje.format(**kwargs)
        else:
            texto = mensaje
        print(f"[DEBUG] {texto}")


# ===========================================================================
# ESTADÍSTICAS DE CONVERSIÓN
# ===========================================================================

class EstadisticasConversion:
    """Almacena métricas de conversión para modo debug."""
    
    def __init__(self):
        self.reset()
    
    def reset(self):
        self.lineas_raw = 0
        self.lineas_limpias = 0
        self.chars_raw = 0
        self.chars_limpios = 0
        self.tablas_detectadas = 0
        self.tablas_convertidas = 0
        self.bloques_duplicados_eliminados = 0
        self.palabras_guionadas_unidas = 0
        self.referencias_normalizadas = 0
        self.archivos_por_tipo = {}
    
    def incrementar_tipo(self, extension: str):
        """Cuenta archivos procesados por tipo."""
        self.archivos_por_tipo[extension] = self.archivos_por_tipo.get(extension, 0) + 1
    
    def mostrar_resumen(self):
        """Imprime resumen de estadísticas en modo debug."""
        if not _DEBUG_MODE:
            return
        
        print("\n" + "─" * 60)
        print("📊 ESTADÍSTICAS DE CONVERSIÓN")
        print("─" * 60)
        
        if self.archivos_por_tipo:
            print("📁 Archivos por tipo:")
            for ext, count in sorted(self.archivos_por_tipo.items()):
                print(f"   {ext:10s} : {count:3d} archivos")
        
        if self.lineas_raw > 0:
            print(f"\n📄 Líneas:")
            print(f"   Raw    : {self.lineas_raw:,}")
            print(f"   Limpias: {self.lineas_limpias:,}")
            print(f"   Reducción: {self.lineas_raw - self.lineas_limpias:,} líneas")
        
        if self.chars_raw > 0:
            print(f"\n📝 Caracteres:")
            print(f"   Raw    : {self.chars_raw:,}")
            print(f"   Limpios: {self.chars_limpios:,}")
            reduccion_pct = (self.chars_raw - self.chars_limpios) / self.chars_raw * 100 if self.chars_raw else 0
            print(f"   Reducción: {self.chars_raw - self.chars_limpios:,} chars ({reduccion_pct:.1f}%)")
        
        if any([self.palabras_guionadas_unidas, self.referencias_normalizadas, self.bloques_duplicados_eliminados]):
            print(f"\n📊 Procesamiento:")
            if self.palabras_guionadas_unidas:
                print(f"   Palabras guionadas unidas      : {self.palabras_guionadas_unidas}")
            if self.referencias_normalizadas:
                print(f"   Referencias normalizadas        : {self.referencias_normalizadas}")
            if self.bloques_duplicados_eliminados:
                print(f"   Bloques duplicados eliminados   : {self.bloques_duplicados_eliminados}")
        
        if self.tablas_detectadas > 0:
            print(f"\n📦 Tablas:")
            print(f"   Detectadas : {self.tablas_detectadas}")
            print(f"   Convertidas: {self.tablas_convertidas}")
            tasa = self.tablas_convertidas / self.tablas_detectadas * 100
            print(f"   Tasa éxito : {tasa:.1f}%")
        
        print("─" * 60 + "\n")


# ===========================================================================
# UTILIDADES
# ===========================================================================

def simple_slug(text: str) -> str:
    """Genera slug seguro para IDs de Docusaurus."""
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("ascii", "ignore").decode("ascii")
    return re.sub(r"[^\w\-]", "-", text.lower()).strip("-")


def get_frontmatter(file_path: Path) -> str:
    """Genera frontmatter YAML para Docusaurus."""
    doc_id = simple_slug(file_path.stem)
    title = file_path.stem.replace("-", " ").replace("_", " ").title()
    return f"""---
id: {doc_id}
title: "{title}"
sidebar_label: "{title}"
---

"""


# ===========================================================================
# PIPELINE DE LIMPIEZA MARKDOWN (Reutilizable)
# ===========================================================================

def unir_palabras_guionadas(texto: str) -> str:
    """Restaura palabras cortadas con guión de fin de línea."""
    contador = len(re.findall(r"([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])-\n([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])", texto))
    _stats.palabras_guionadas_unidas += contador
    debug_log("Uniendo {n} palabras guionadas", n=contador)
    return re.sub(r"([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])-\n([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])", r"\1\2", texto)


def limpiar_indice(texto: str) -> str:
    """Limpia entradas de índices (puntos suspensivos + números de página)."""
    RE_TIENE_DOTS = re.compile(r"\.{3,}\s*\d")
    RE_QUITAR_DOTS_NUM = re.compile(r"\s*\.{3,}\s*\d+")
    
    lineas = texto.split("\n")
    resultado = []
    lineas_eliminadas = 0
    
    for linea in lineas:
        if not RE_TIENE_DOTS.search(linea):
            resultado.append(linea)
            continue
        
        limpia = RE_QUITAR_DOTS_NUM.sub("", linea).strip()
        if limpia:
            resultado.append(limpia)
        else:
            lineas_eliminadas += 1
    
    debug_log("Índice: {n} líneas de numeración eliminadas", n=lineas_eliminadas)
    return "\n".join(resultado)


def limpiar_encabezados_pies(texto: str) -> str:
    """Elimina encabezados y pies de página repetitivos."""
    patron_num_norma = r"^\s*\d+\s*[-–—]\s*[A-Za-z]{2,}[\w/\-\s.:]+\d{4,}\s*$"
    patron_norma_num = r"^[A-Za-z]{2,}[\w/\-\s.:]+\d{4,}\s*[-–—]\s*\d+\s*$"
    patron_pagina = r"^\s*\d+\s*$"
    
    lineas = texto.split("\n")
    lineas_eliminadas = 0
    resultado = []
    
    for linea in lineas:
        if (re.match(patron_num_norma, linea) or
            re.match(patron_norma_num, linea) or
            re.match(patron_pagina, linea)):
            lineas_eliminadas += 1
            continue
        resultado.append(linea)
    
    debug_log("Encabezados/pies eliminados: {n} líneas", n=lineas_eliminadas)
    return "\n".join(resultado)


def fusionar_parrafos_rotos(texto: str) -> str:
    """Fusiona líneas cortadas arbitrariamente por la conversión."""
    FIN_FRASE = (".", ":", ";", "!", "?", ")")
    MARCADORES = ("#", "-", "*", "|", ">", "```")
    
    lineas = texto.split("\n")
    lineas_fusionadas = 0
    resultado = []
    i = 0
    
    while i < len(lineas):
        linea = lineas[i].rstrip()
        
        if not linea:
            resultado.append("")
            i += 1
            continue
        
        # No fusionar líneas protegidas
        if linea.startswith(MARCADORES) or re.match(r"^\d+(?:\.\d+)*\s", linea):
            resultado.append(linea)
            i += 1
            continue
        
        # No fusionar si termina en puntuación
        if linea.endswith(FIN_FRASE):
            resultado.append(linea)
            i += 1
            continue
        
        # Intentar fusionar con siguiente
        if i + 1 < len(lineas):
            sig = lineas[i + 1].lstrip()
            if not sig or sig.startswith(MARCADORES):
                resultado.append(linea)
                i += 1
            else:
                resultado.append(linea + " " + sig)
                lineas_fusionadas += 1
                i += 2
        else:
            resultado.append(linea)
            i += 1
    
    debug_log("Párrafos fusionados: {n} líneas unidas", n=lineas_fusionadas)
    return "\n".join(resultado)


def normalizar_referencias(texto: str) -> str:
    """Limpia referencias cruzadas y citas normativas."""
    contador = 0
    
    # Referencias partidas
    nuevo, n1 = re.subn(
        r"(\b(?:véase|see|clause|cláusula|apartado|section)\b)\n(\d[\d\.]*)",
        r"\1 \2",
        texto,
        flags=re.IGNORECASE,
    )
    contador += n1
    
    # Espacios en citas
    nuevo, n2 = re.subn(
        r"\[([A-Z]{2,}[/\-\w\s]+?)\s*:\s*(\d{4})\]",
        r"[\1:\2]",
        nuevo,
    )
    contador += n2
    
    _stats.referencias_normalizadas += contador
    debug_log("Referencias normalizadas: {n}", n=contador)
    return nuevo


def arreglar_html_mdx(texto: str) -> str:
    """Convierte etiquetas HTML a forma JSX auto-cerrada."""
    texto = re.sub(r"<br\s*>", "<br />", texto, flags=re.IGNORECASE)
    texto = re.sub(r"<hr\s*>", "<hr />", texto, flags=re.IGNORECASE)
    texto = re.sub(r"<img([^>]+)(?<!/)>", r"<img\1 />", texto, flags=re.IGNORECASE)
    return texto


def limpiar_espacios_multiples(texto: str) -> str:
    """Normaliza espacios múltiples y líneas en blanco."""
    texto = re.sub(r" {2,}", " ", texto)
    texto = re.sub(r"\t", "    ", texto)
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    return texto


def limpiar_markdown_completo(md_texto: str, verbose: bool = True) -> str:
    """
    Pipeline completo de limpieza Markdown.
    Reutilizable para cualquier fuente (PDF, DOCX, etc.)
    """
    _stats.chars_raw += len(md_texto)
    _stats.lineas_raw += md_texto.count("\n") + 1
    
    debug_log("=" * 60)
    debug_log("INICIANDO PIPELINE DE LIMPIEZA")
    debug_log("=" * 60)
    
    pasos = [
        ("🔗 Uniendo palabras guionadas", unir_palabras_guionadas),
        ("🧹 Limpiando índice", limpiar_indice),
        ("🧹 Eliminando encabezados/pies", limpiar_encabezados_pies),
        ("🔗 Normalizando referencias", normalizar_referencias),
        ("🧹 Fusionando párrafos rotos", fusionar_parrafos_rotos),
        ("🧹 Limpiando espacios", limpiar_espacios_multiples),
        ("🧹 Arreglando HTML/MDX", arreglar_html_mdx),
    ]
    
    texto = md_texto
    for nombre, fn in pasos:
        if verbose:
            print(f"      {nombre}...")
        texto = fn(texto)
    
    _stats.chars_limpios += len(texto)
    _stats.lineas_limpias += texto.count("\n") + 1
    
    return texto.strip()


# ===========================================================================
# HANDLERS POR TIPO DE ARCHIVO
# ===========================================================================

def handle_pdf(file_path: Path) -> Optional[str]:
    """
    Convierte PDF a Markdown usando pymupdf4llm.
    Retorna: Markdown raw o None si falla.
    """
    if not PYMUPDF_DISPONIBLE:
        print(f"   ⚠️  pymupdf4llm no disponible - Saltando PDF")
        return None
    
    try:
        debug_log("Extrayendo PDF con pymupdf4llm...")
        md_raw = pymupdf4llm.to_markdown(str(file_path))
        debug_log("PDF extraído: {n} caracteres", n=len(md_raw))
        return md_raw
    except Exception as e:
        print(f"   ❌ Error extrayendo PDF: {e}")
        return None


def handle_docx(file_path: Path) -> Optional[str]:
    """
    Convierte DOCX a Markdown.
    Retorna: Markdown raw o None si falla.
    """
    if not DOCX_DISPONIBLE:
        print(f"   ⚠️  python-docx no disponible - Saltando DOCX")
        return None
    
    try:
        debug_log("Extrayendo DOCX con python-docx...")
        doc = Document(file_path)
        
        markdown_lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                markdown_lines.append("")
                continue
            
            # Detectar estilos de título
            if para.style.name.startswith('Heading'):
                level = para.style.name.replace('Heading ', '')
                try:
                    level_num = int(level)
                    markdown_lines.append(f"{'#' * level_num} {text}")
                except:
                    markdown_lines.append(text)
            else:
                markdown_lines.append(text)
        
        # Procesar tablas
        for table in doc.tables:
            markdown_lines.append("")
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                markdown_lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    markdown_lines.append("|" + "---|" * len(cells))
            markdown_lines.append("")
        
        md_raw = "\n".join(markdown_lines)
        debug_log("DOCX extraído: {n} caracteres", n=len(md_raw))
        return md_raw
    
    except Exception as e:
        print(f"   ❌ Error extrayendo DOCX: {e}")
        return None


def handle_excel(file_path: Path) -> Optional[str]:
    """
    Convierte Excel (XLSX, XLS, XLSM) a Markdown.
    Cada hoja se convierte en una tabla Markdown.
    """
    if not PANDAS_DISPONIBLE:
        print(f"   ⚠️  pandas no disponible - Saltando Excel")
        return None
    
    try:
        debug_log("Extrayendo Excel con pandas...")
        
        # Leer todas las hojas
        excel_file = pd.ExcelFile(file_path)
        markdown_lines = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Título de hoja
            markdown_lines.append(f"\n## {sheet_name}\n")
            
            # Convertir a tabla Markdown
            if not df.empty:
                # Encabezados
                headers = df.columns.tolist()
                markdown_lines.append("| " + " | ".join(str(h) for h in headers) + " |")
                markdown_lines.append("|" + "---|" * len(headers))
                
                # Filas
                for _, row in df.iterrows():
                    cells = [str(val) if pd.notna(val) else "" for val in row]
                    markdown_lines.append("| " + " | ".join(cells) + " |")
                
                markdown_lines.append("")
                _stats.tablas_detectadas += 1
                _stats.tablas_convertidas += 1
        
        md_raw = "\n".join(markdown_lines)
        debug_log("Excel extraído: {n} hojas, {chars} caracteres", 
                  n=len(excel_file.sheet_names), chars=len(md_raw))
        return md_raw
    
    except Exception as e:
        print(f"   ❌ Error extrayendo Excel: {e}")
        return None


def handle_pptx(file_path: Path) -> Optional[str]:
    """
    Convierte PowerPoint (PPTX) a Markdown.
    Cada diapositiva se convierte en una sección.
    """
    if not PPTX_DISPONIBLE:
        print(f"   ⚠️  python-pptx no disponible - Saltando PPTX")
        return None
    
    try:
        debug_log("Extrayendo PPTX con python-pptx...")
        prs = Presentation(file_path)
        markdown_lines = []
        
        for i, slide in enumerate(prs.slides, 1):
            markdown_lines.append(f"\n## Diapositiva {i}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    markdown_lines.append(shape.text.strip())
                    markdown_lines.append("")
                
                # Procesar tablas si existen
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        markdown_lines.append("| " + " | ".join(cells) + " |")
                        if row_idx == 0:
                            markdown_lines.append("|" + "---|" * len(cells))
                    markdown_lines.append("")
        
        md_raw = "\n".join(markdown_lines)
        debug_log("PPTX extraído: {n} diapositivas, {chars} caracteres", 
                  n=len(prs.slides), chars=len(md_raw))
        return md_raw
    
    except Exception as e:
        print(f"   ❌ Error extrayendo PPTX: {e}")
        return None


def handle_txt(file_path: Path) -> Optional[str]:
    """
    Lee archivo TXT y lo convierte a Markdown (mínimo procesamiento).
    """
    try:
        debug_log("Leyendo archivo TXT...")
        md_raw = file_path.read_text(encoding='utf-8')
        debug_log("TXT leído: {n} caracteres", n=len(md_raw))
        return md_raw
    except UnicodeDecodeError:
        try:
            # Intentar con latin-1
            md_raw = file_path.read_text(encoding='latin-1')
            debug_log("TXT leído (latin-1): {n} caracteres", n=len(md_raw))
            return md_raw
        except Exception as e:
            print(f"   ❌ Error leyendo TXT: {e}")
            return None
    except Exception as e:
        print(f"   ❌ Error leyendo TXT: {e}")
        return None


def handle_png(file_path: Path) -> Optional[str]:
    """
    Extrae texto de imagen PNG usando OCR (Tesseract).
    Requiere: pytesseract y Tesseract instalado en el sistema.
    """
    if not OCR_DISPONIBLE:
        print(f"   ⚠️  pytesseract no disponible - Saltando PNG")
        return None
    
    try:
        debug_log("Extrayendo texto de PNG con OCR...")
        image = Image.open(file_path)
        md_raw = pytesseract.image_to_string(image, lang='spa')
        debug_log("PNG extraído: {n} caracteres", n=len(md_raw))
        return md_raw
    except Exception as e:
        print(f"   ❌ Error procesando PNG: {e}")
        return None


def handle_zip(file_path: Path) -> Optional[str]:
    """
    Extrae contenido de archivo ZIP.
    Por ahora, solo lista los archivos contenidos.
    """
    try:
        import zipfile
        debug_log("Listando contenido de ZIP...")
        
        markdown_lines = [f"# Contenido de {file_path.name}\n"]
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            file_list = zip_ref.namelist()
            markdown_lines.append(f"Total de archivos: {len(file_list)}\n")
            
            for filename in file_list:
                markdown_lines.append(f"- {filename}")
        
        md_raw = "\n".join(markdown_lines)
        debug_log("ZIP listado: {n} archivos", n=len(file_list))
        return md_raw
    
    except Exception as e:
        print(f"   ❌ Error procesando ZIP: {e}")
        return None


# ===========================================================================
# DISPATCHER DE HANDLERS
# ===========================================================================

# Mapeo de extensiones a funciones handler
HANDLERS = {
    '.pdf': handle_pdf,
    '.docx': handle_docx,
    '.xlsx': handle_excel,
    '.xls': handle_excel,
    '.xlsm': handle_excel,
    '.pptx': handle_pptx,
    '.txt': handle_txt,
    '.png': handle_png,
    '.zip': handle_zip,
}

# Extensiones que se ignoran con warning
EXTENSIONES_IGNORADAS = {'.exe', '.url'}


def convertir_archivo(file_path: Path, output_path: Path, verbose: bool = True) -> bool:
    """
    Convierte un archivo al formato Markdown según su extensión.
    
    Returns:
        True si la conversión fue exitosa, False en caso contrario.
    """
    extension = file_path.suffix.lower()
    
    # Verificar si es una extensión ignorada
    if extension in EXTENSIONES_IGNORADAS:
        if verbose:
            print(f"\n⚠️  Ignorando: {file_path.name} (extensión {extension} no soportada)")
        return False
    
    # Verificar si tenemos handler para esta extensión
    if extension not in HANDLERS:
        if verbose:
            print(f"\n⚠️  Formato desconocido: {file_path.name} (extensión {extension})")
        return False
    
    try:
        t_inicio = time.time()
        
        if verbose:
            print(f"\n📖 Procesando: {file_path.name}")
            print(f"   📦 Tipo: {extension.upper()[1:]}")
        
        # Obtener handler correspondiente
        handler = HANDLERS[extension]
        
        # Extraer contenido
        if verbose:
            print(f"   ⚙️  Extrayendo contenido...")
        
        md_raw = handler(file_path)
        
        if md_raw is None:
            return False
        
        if verbose:
            print(f"   📊 Contenido extraído: {len(md_raw):,} caracteres")
            print(f"   🔧 Aplicando limpieza Markdown...")
        
        # Aplicar pipeline de limpieza
        md_limpio = limpiar_markdown_completo(md_raw, verbose=verbose)
        
        # Añadir frontmatter
        contenido_final = get_frontmatter(file_path) + md_limpio
        
        # Guardar archivo
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(contenido_final, encoding='utf-8')
        
        # Estadísticas
        _stats.incrementar_tipo(extension)
        
        t_total = time.time() - t_inicio
        
        if verbose:
            reduccion = len(md_raw) - len(md_limpio)
            pct = reduccion / len(md_raw) * 100 if md_raw else 0
            print(f"   ✅ Guardado: {output_path}")
            print(f"   📉 Reducción: {reduccion:,} chars ({pct:.1f}%)")
        
        debug_log("⏱  Tiempo de conversión: {t:.2f}s", t=t_total)
        
        return True
    
    except Exception as e:
        print(f"   ❌ ERROR en {file_path.name}: {e}")
        if _DEBUG_MODE:
            import traceback
            traceback.print_exc()
        return False


# ===========================================================================
# PROCESAMIENTO MASIVO
# ===========================================================================

def procesar_carpeta(
    input_dir: Path,
    output_dir: Path,
    verbose: bool = True,
    workers: int = 1,
    extensiones_filtro: Optional[list] = None,
):
    """
    Procesa todos los archivos de una carpeta.
    
    Args:
        input_dir: Carpeta de origen
        output_dir: Carpeta de destino
        verbose: Mostrar mensajes detallados
        workers: Número de hilos paralelos
        extensiones_filtro: Lista de extensiones a procesar (ej: ['.pdf', '.docx'])
                            Si es None, procesa todas las soportadas.
    """
    separador = "=" * 70
    
    if verbose:
        print(f"\n{separador}")
        print(f"🔄 INICIANDO CONVERSIÓN MULTI-FORMATO — v4")
        print(f"{separador}")
        print(f"📁 Origen : {input_dir}")
        print(f"📂 Destino: {output_dir}")
        print(f"⚡ Workers: {workers}")
        if _DEBUG_MODE:
            print(f"🐛 Modo DEBUG: ACTIVO")
    
    # Buscar todos los archivos soportados
    archivos = []
    extensiones_buscar = extensiones_filtro or list(HANDLERS.keys())
    
    for ext in extensiones_buscar:
        archivos.extend(input_dir.rglob(f"*{ext}"))
    
    if not archivos:
        print(f"❌ No se encontraron archivos con extensiones soportadas.")
        print(f"   Extensiones buscadas: {', '.join(extensiones_buscar)}")
        return
    
    print(f"📊 Total archivos encontrados: {len(archivos)}\n")
    t_inicio = time.time()
    
    exitos = 0
    errores = 0
    ignorados = 0
    
    def _tarea(file_path: Path):
        rel_path = file_path.relative_to(input_dir)
        out_path = output_dir / rel_path.with_suffix('.md')
        return convertir_archivo(file_path, out_path, verbose=(workers == 1 and verbose))
    
    if workers == 1:
        # Modo secuencial
        for file_path in archivos:
            if file_path.suffix.lower() in EXTENSIONES_IGNORADAS:
                ignorados += 1
                continue
            
            if _tarea(file_path):
                exitos += 1
            else:
                errores += 1
    else:
        # Modo paralelo
        with ThreadPoolExecutor(max_workers=workers) as executor:
            futuros = {}
            for file_path in archivos:
                if file_path.suffix.lower() in EXTENSIONES_IGNORADAS:
                    ignorados += 1
                    continue
                futuros[executor.submit(_tarea, file_path)] = file_path
            
            procesados = 0
            for futuro in as_completed(futuros):
                procesados += 1
                nombre = futuros[futuro].name
                ok = futuro.result()
                estado = "✅" if ok else "❌"
                print(f"  [{procesados}/{len(archivos) - ignorados}] {estado} {nombre}")
                if ok:
                    exitos += 1
                else:
                    errores += 1
    
    t_total = time.time() - t_inicio
    velocidad = (exitos + errores) / (t_total / 60) if t_total > 0 else 0
    
    # Mostrar estadísticas de conversión
    _stats.mostrar_resumen()
    
    print(f"\n{separador}")
    print(f"📊 RESUMEN FINAL")
    print(f"{separador}")
    print(f"✅ Conversiones exitosas : {exitos}")
    print(f"❌ Errores              : {errores}")
    print(f"⏭️  Archivos ignorados    : {ignorados}")
    print(f"📁 Total procesados     : {exitos + errores}")
    print(f"⏱  Tiempo total         : {t_total:.1f}s")
    print(f"⚡ Velocidad            : {velocidad:.1f} archivos/min")
    print(f"📂 Archivos en          : {output_dir}")
    print(f"{separador}\n")


# ===========================================================================
# CLI
# ===========================================================================

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Conversor multi-formato a Markdown para Docusaurus (GRATIS) — v4",
        epilog=(
            "Formatos soportados: PDF, DOCX, XLSX, XLS, XLSM, PPTX, TXT, PNG, ZIP\n"
            "\n"
            "Ejemplos:\n"
            "  python convert_multi_format_v4.py\n"
            "  python convert_multi_format_v4.py --input mis-archivos/ --output docs/\n"
            "  python convert_multi_format_v4.py --workers 4\n"
            "  python convert_multi_format_v4.py --debug\n"
            "  python convert_multi_format_v4.py --only-pdf\n"
            "  python convert_multi_format_v4.py --only-excel\n"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    
    parser.add_argument("--input", default="pdf", 
                        help="Carpeta origen (default: pdf/)")
    parser.add_argument("--output", default="docs", 
                        help="Carpeta destino (default: docs/)")
    parser.add_argument("--workers", default=1, type=int,
                        help="Hilos paralelos (default: 1)")
    parser.add_argument("--debug", action="store_true",
                        help="Modo debug con estadísticas detalladas")
    parser.add_argument("--quiet", action="store_true",
                        help="Modo silencioso")
    
    # Filtros por tipo de archivo
    parser.add_argument("--only-pdf", action="store_true",
                        help="Solo procesar archivos PDF")
    parser.add_argument("--only-docx", action="store_true",
                        help="Solo procesar archivos DOCX")
    parser.add_argument("--only-excel", action="store_true",
                        help="Solo procesar archivos Excel (XLSX, XLS, XLSM)")
    parser.add_argument("--only-pptx", action="store_true",
                        help="Solo procesar archivos PowerPoint")
    
    args = parser.parse_args()
    
    # Activar modo debug
    if args.debug:
        set_debug_mode(True)
    
    # Inicializar estadísticas globales
    _stats = EstadisticasConversion()
    
    # Determinar filtro de extensiones
    extensiones_filtro = None
    if args.only_pdf:
        extensiones_filtro = ['.pdf']
    elif args.only_docx:
        extensiones_filtro = ['.docx']
    elif args.only_excel:
        extensiones_filtro = ['.xlsx', '.xls', '.xlsm']
    elif args.only_pptx:
        extensiones_filtro = ['.pptx']
    
    # Validar carpetas
    input_path = Path(args.input)
    output_path = Path(args.output)
    
    if not input_path.exists():
        print(f"⚠️  Carpeta '{input_path}' no existe. Creándola...")
        input_path.mkdir(parents=True)
        print(f"   ℹ️  Coloca tus archivos en '{input_path}' y ejecuta de nuevo.")
        raise SystemExit(0)
    
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Procesar
    procesar_carpeta(
        input_path,
        output_path,
        verbose=not args.quiet,
        workers=args.workers,
        extensiones_filtro=extensiones_filtro,
    )

