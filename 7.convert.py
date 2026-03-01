# convert_multi_format_v6.py
# Conversor modular multi-formato para Docusaurus — v6 (STRICT MODE)
# Soporta: PDF, DOCX, XLSX, XLS, XLSM, PPTX, TXT, PNG, ZIP
# 100% GRATIS - Sin APIs de pago
#
# NOVEDADES v6:
#   - 🛡️ PROTECCIÓN DE TABLAS: Aísla tablas del pipeline de limpieza para evitar destrozos.
#   - 🧹 LIMPIEZA DE FIRMAS: Elimina metadatos de firmas digitales de PDFs.
#   - ⚖️ JUSTIFICACIÓN ROBUSTA: Solo aplica a párrafos continuos (excluye estructura).
#   - ⚙️ MODO CONSERVADOR: --conservar-tablas para documentos críticos.
#   - 🐛 MDX SAFE: Escapado automático de < y { para evitar errores de compilación.

import re
import time
import unicodedata
import logging
import tempfile
import shutil
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Optional, Tuple, List, Dict
from datetime import datetime
import warnings

# Dependencias según disponibilidad
try:
    import fitz  # PyMuPDF
    import pymupdf4llm
    PYMUPDF_DISPONIBLE = True
except ImportError:
    PYMUPDF_DISPONIBLE = False
    warnings.warn("pymupdf4llm no disponible - Los PDFs no se podrán procesar")

try:
    from docx import Document
    DOCX_DISPONIBLE = True
except ImportError:
    DOCX_DISPONIBLE = False
    warnings.warn("python-docx no disponible - Los DOCX no se podrán procesar")

try:
    import pandas as pd
    PANDAS_DISPONIBLE = True
except ImportError:
    PANDAS_DISPONIBLE = False
    warnings.warn("pandas no disponible - Los archivos Excel no se podrán procesar")

try:
    from pptx import Presentation
    PPTX_DISPONIBLE = True
except ImportError:
    PPTX_DISPONIBLE = False
    warnings.warn("python-pptx no disponible - Los PPTX no se podrán procesar")

try:
    from PIL import Image
    import pytesseract
    OCR_DISPONIBLE = True
except ImportError:
    OCR_DISPONIBLE = False
    warnings.warn("PIL/pytesseract no disponibles - OCR no estará disponible")


# ===========================================================================
# CONFIGURACIÓN GLOBAL
# ===========================================================================

_DEBUG_MODE = False
_OCR_MODE = False
_STRICT_TABLES = False
_stats = None
_logger = None


def set_modes(debug: bool = False, ocr: bool = False, strict: bool = False):
    global _DEBUG_MODE, _OCR_MODE, _STRICT_TABLES
    _DEBUG_MODE = debug
    _OCR_MODE = ocr
    _STRICT_TABLES = strict

def debug_log(mensaje: str, **kwargs):
    """Imprime mensajes solo si el modo debug está activo."""
    if _DEBUG_MODE:
        if kwargs:
            texto = mensaje.format(**kwargs)
        else:
            texto = mensaje
        print(f"[DEBUG] {texto}")


def setup_logger(output_dir: Path):
    """Configura logger para registrar errores en archivo."""
    global _logger
    log_file = output_dir / "errors.log"
    
    _logger = logging.getLogger('converter')
    _logger.setLevel(logging.ERROR)
    
    # Handler para archivo
    fh = logging.FileHandler(log_file, encoding='utf-8')
    fh.setLevel(logging.ERROR)
    
    # Formato
    formatter = logging.Formatter(
        '%(asctime)s - %(levelname)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    fh.setFormatter(formatter)
    
    _logger.addHandler(fh)
    
    return log_file


# ===========================================================================
# ESTADÍSTICAS DE CONVERSIÓN
# ===========================================================================

class EstadisticasConversion:
    """Almacena métricas de conversión para modo debug."""
    
    def __init__(self):
        self.reset()
    
    def reset(self):
        self.lineas_raw = 0
        self.lineas_limpias = 0
        self.chars_raw = 0
        self.chars_limpios = 0
        self.tablas_detectadas = 0
        self.tablas_convertidas = 0
        self.bloques_duplicados_eliminados = 0
        self.palabras_guionadas_unidas = 0
        self.referencias_normalizadas = 0
        self.firmas_digitales_limpiadas = 0
        self.archivos_por_tipo = {}
        self.archivos_con_ocr = 0
        self.pdfs_escaneados = 0
        self.zips_procesados = 0
        self.archivos_en_zips = 0
        self.tiempo_total = 0.0
    
    def incrementar_tipo(self, extension: str):
        """Cuenta archivos procesados por tipo."""
        self.archivos_por_tipo[extension] = self.archivos_por_tipo.get(extension, 0) + 1
    
    def mostrar_resumen(self):
        """Imprime resumen de estadísticas en modo debug."""
        if not _DEBUG_MODE:
            return
        
        print("\n" + "─" * 60)
        print("📊 ESTADÍSTICAS DE CONVERSIÓN (v6)")
        print("─" * 60)
        
        if self.archivos_por_tipo:
            print("📁 Archivos por tipo:")
            for ext, count in sorted(self.archivos_por_tipo.items()):
                print(f"   {ext:10s} : {count:3d} archivos")
        
        if self.pdfs_escaneados > 0:
            print(f"\n📄 PDFs escaneados detectados: {self.pdfs_escaneados}")
        
        if self.archivos_con_ocr > 0:
            print(f"🔍 Archivos procesados con OCR: {self.archivos_con_ocr}")
        
        if self.zips_procesados > 0:
            print(f"\n📦 Archivos ZIP:")
            print(f"   ZIPs procesados     : {self.zips_procesados}")
            print(f"   Archivos extraídos  : {self.archivos_en_zips}")
        
        if self.tablas_detectadas > 0:
            print(f"\n📊 Tablas:")
            print(f"   Detectadas  : {self.tablas_detectadas}")
            print(f"   Convertidas : {self.tablas_convertidas}")

        if self.firmas_digitales_limpiadas > 0:
            print(f"✍️  Firmas digitales limpiadas: {self.firmas_digitales_limpiadas}")

        if self.lineas_raw > 0:
            print(f"\n📄 Líneas:")
            print(f"   Raw    : {self.lineas_raw:,}")
            print(f"   Limpias: {self.lineas_limpias:,}")
        
        if self.chars_raw > 0:
            print(f"\n📝 Caracteres:")
            print(f"   Raw    : {self.chars_raw:,}")
            print(f"   Limpios: {self.chars_limpios:,}")
            reduccion_pct = (self.chars_raw - self.chars_limpios) / self.chars_raw * 100 if self.chars_raw else 0
            print(f"   Reducción: {self.chars_raw - self.chars_limpios:,} chars ({reduccion_pct:.1f}%)")

        if self.tiempo_total > 0:
             print(f"\n⏱  Tiempo total de ejecución: {self.tiempo_total:.2f}s")
        
        print("─" * 60 + "\n")


# ===========================================================================
# UTILIDADES
# ===========================================================================

def normalizar_nombre_archivo(nombre: str) -> str:
    """Normaliza nombre de archivo para evitar problemas de caracteres especiales."""
    ruta = Path(nombre)
    stem = ruta.stem
    ext = ruta.suffix.lower()
    
    stem_normalizado = unicodedata.normalize("NFKD", stem)
    stem_normalizado = stem_normalizado.encode("ascii", "ignore").decode("ascii")
    # Mantener más contexto pero eliminar peligrosos
    stem_normalizado = re.sub(r"[^\w\s\-]", "", stem_normalizado).strip()
    stem_normalizado = re.sub(r"[\s_]+", "-", stem_normalizado.lower())
    
    if not stem_normalizado:
        stem_normalizado = "documento"
    
    return stem_normalizado + ext


def simple_slug(text: str) -> str:
    """Genera slug seguro para IDs de Docusaurus."""
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("ascii", "ignore").decode("ascii")
    return re.sub(r"[^\w\-]", "-", text.lower()).strip("-")


def get_frontmatter(file_path: Path) -> str:
    """Genera frontmatter YAML para Docusaurus."""
    doc_id = simple_slug(file_path.stem)
    title = file_path.stem.replace("-", " ").replace("_", " ").title()
    return f"""---
id: {doc_id}
title: "{title}"
sidebar_label: "{title}"
---

"""


# ===========================================================================
# PROTECCIÓN DE TABLAS (Placeholder Mechanism)
# ===========================================================================

def aislar_tablas(texto: str) -> Tuple[str, List[str]]:
    """Extrae tablas de Markdown y las reemplaza por placeholders."""
    debug_log("Aislando tablas para protección...")
    tablas = []
    lineas = texto.split('\n')
    resultado = []
    
    i = 0
    while i < len(lineas):
        stripped = lineas[i].strip()
        # Una tabla en MD suele empezar con |
        if stripped.startswith('|'):
            bloque_tabla = []
            while i < len(lineas) and lineas[i].strip().startswith('|'):
                bloque_tabla.append(lineas[i])
                i += 1
            
            contenido_tabla = '\n'.join(bloque_tabla)
            placeholder = f"<!-- TABLE_PLACEHOLDER_{len(tablas)} -->"
            tablas.append(contenido_tabla)
            resultado.append(placeholder)
        else:
            resultado.append(lineas[i])
            i += 1
            
    return '\n'.join(resultado), tablas


def reinsertar_tablas(texto: str, tablas: List[str]) -> str:
    """Restaura las tablas en sus placeholders correspondientes."""
    debug_log("Reinsertando tablas protegidas...")
    for idx, tabla in enumerate(tablas):
        placeholder = f"<!-- TABLE_PLACEHOLDER_{idx} -->"
        texto = texto.replace(placeholder, f"\n{tabla}\n")
    return texto


# ===========================================================================
# PIPELINE DE LIMPIEZA MARKDOWN
# ===========================================================================

def limpiar_firmas(texto: str) -> str:
    """Elimina o simplifica texto de firmas digitales."""
    original = len(texto)
    
    # Patrón común de firmas PDF
    texto = re.sub(r'`Firmado por[^`]+`', '[FIRMA DIGITAL]', texto, flags=re.IGNORECASE)
    # Patrón de asteriscos y números de certificados
    texto = re.sub(r'\*{2,}\d+\*{2,}', '', texto)
    # Líneas colgantes de fechas de firma
    texto = re.sub(r'el día \d{2}/\d{2}/\d{4}.*$', '', texto, flags=re.MULTILINE)
    
    if len(texto) < original:
        _stats.firmas_digitales_limpiadas += 1
        
    return texto


def unir_palabras_guionadas(texto: str) -> str:
    """Restaura palabras cortadas con guión de fin de línea."""
    return re.sub(r"([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])-\n([a-záéíóúüñA-ZÁÉÍÓÚÜÑ])", r"\1\2", texto)


def limpiar_indice(texto: str) -> str:
    """Limpia entradas de índices con puntos extensivos."""
    RE_QUITAR_DOTS_NUM = re.compile(r"\s*\.{3,}\s*\d+")
    lineas = texto.split("\n")
    resultado = [RE_QUITAR_DOTS_NUM.sub("", l).strip() for l in lineas]
    return "\n".join([l for l in resultado if l])


def limpiar_encabezados_pies(texto: str) -> str:
    """Elimina encabezados y pies de página numéricos o repetitivos."""
    patrones = [
        r"^\s*\d+\s*[-–—]\s*[A-Za-z]{2,}[\w/\-\s.:]+\d{4,}\s*$",
        r"^[A-Za-z]{2,}[\w/\-\s.:]+\d{4,}\s*[-–—]\s*\d+\s*$",
        r"^\s*\d+\s*$"
    ]
    lineas = texto.split("\n")
    resultado = [l for l in lineas if not any(re.match(p, l) for p in patrones)]
    return "\n".join(resultado)


def fusionar_parrafos_rotos(texto: str) -> str:
    """Fusiona líneas que no terminan en puntuación clara."""
    if _STRICT_TABLES: # En modo estricto, somos menos agresivos
        return texto
        
    FIN_FRASE = (".", ":", ";", "!", "?", ")", "]", '"', "'")
    MARCADORES = ("#", "-", "*", "|", ">", "```", "<!--")
    
    lineas = texto.split("\n")
    resultado = []
    i = 0
    
    while i < len(lineas):
        linea = lineas[i].rstrip()
        if not linea or linea.startswith(MARCADORES) or re.match(r"^\d+(?:\.\d+)*\s", linea) or linea.endswith(FIN_FRASE):
            resultado.append(linea)
            i += 1
            continue
        
        if i + 1 < len(lineas):
            sig = lineas[i + 1].lstrip()
            if not sig or sig.startswith(MARCADORES):
                resultado.append(linea)
                i += 1
            else:
                resultado.append(linea + " " + sig)
                i += 2
        else:
            resultado.append(linea)
            i += 1
            
    return "\n".join(resultado)


def arreglar_html_mdx(texto: str) -> str:
    """Asegura compatibilidad MDX (JSX tags, escapes)."""
    # 1. JSX Compatibility
    texto = re.sub(r"<br\s*>", "<br />", texto, flags=re.IGNORECASE)
    texto = re.sub(r"<hr\s*>", "<hr />", texto, flags=re.IGNORECASE)
    texto = re.sub(r"<img([^>]+)(?<!/)>", r"<img\1 />", texto, flags=re.IGNORECASE)
    
    # 2. Escapar caracteres que MDX interpreta como JSX
    # Escapar < si no es tag válido
    texto = re.sub(r"<(?!/?(?:p|br|hr|img|div|span|strong|em|b|i|u|h\d)\b)", "&lt;", texto)
    # Escapar llaves (conflictos con expresiones JS)
    texto = texto.replace("{", "&#123;").replace("}", "&#125;")
    
    return texto


def justificar_texto(texto: str) -> str:
    """Justifica BLOQUES de texto continuo usando regex."""
    debug_log("Iniciando justificación (v6 engine)...")
    
    # Marcadores estructurales que rompen un párrafo
    # Incluimos los placeholders de tabla para no justificarlos
    PATRON_BREAK = r"^(?:#|\-|\*|\d+\.|\||```|<img|<p|<div|<br|<hr|&lt;|&#123;|--|id:|title:|sidebar_label:|<!--)"
    
    def callback(match):
        bloque = match.group(0)
        lines = bloque.strip().split('\n')
        
        # Si el bloque empieza con estructura o es muy corto, no justificar
        if any(re.match(PATRON_BREAK, l.strip(), re.IGNORECASE) for l in lines) or len(bloque.strip()) < 60:
            return bloque
            
        # Párrafo real -> Envolver en div justificando
        return f'\n<div style={{{{textAlign: "justify"}}}}>\n\n{bloque.strip()}\n\n</div>\n'

    # Captura bloques delimitados por doble salto de línea
    regex_bloques = r"((?:(?!\n\n).)+)"
    return re.sub(regex_bloques, callback, texto, flags=re.DOTALL)


def limpiar_markdown_completo(md_texto: str, verbose: bool = True) -> str:
    """Pipeline v6 con protección de tablas y limpieza de firmas."""
    _stats.chars_raw += len(md_texto)
    _stats.lineas_raw += md_texto.count("\n") + 1
    
    # 1. Limpieza de firmas (antes de aislar para capturar ruido)
    texto = limpiar_firmas(md_texto)
    
    # 2. Arreglar HTML/MDX base
    texto = arreglar_html_mdx(texto)
    
    # 3. AISLAR TABLAS
    texto_limpiable, tablas = aislar_tablas(texto)
    
    # 4. Pipeline de limpieza en el texto restante
    pasos = [
        ("🔗 Uniendo palabras guionadas", unir_palabras_guionadas),
        ("🧹 Limpiando índice", limpiar_indice),
        ("🧹 Eliminando encabezados/pies", limpiar_encabezados_pies),
        ("🧹 Fusionando párrafos", fusionar_parrafos_rotos),
        ("🧹 Limpiando espacios", lambda t: re.sub(r" {2,}", " ", t)),
    ]
    
    for nombre, fn in pasos:
        if verbose: print(f"      {nombre}...")
        texto_limpiable = fn(texto_limpiable)
    
    # 5. Justificar párrafos
    texto_justificado = justificar_texto(texto_limpiable)
    
    # 6. REINSERTAR TABLAS
    texto_final = reinsertar_tablas(texto_justificado, tablas)
    
    # 7. Toques finales
    texto_final = re.sub(r"\n{3,}", "\n\n", texto_final).strip()
    
    _stats.chars_limpios += len(texto_final)
    _stats.lineas_limpias += texto_final.count("\n") + 1
    
    return texto_final


# ===========================================================================
# HANDLERS (Iguales a v5 pero adaptados a v6 si es necesario)
# ===========================================================================

def handle_pdf(file_path: Path) -> Optional[str]:
    if not PYMUPDF_DISPONIBLE: return None
    try:
        # Lógica de detección de escaneado simplificada
        doc = fitz.open(file_path)
        texto_ext = "".join([p.get_text() for p in doc])
        es_escaneado = (len(texto_ext) / len(doc)) < 500 if len(doc) > 0 else True
        doc.close()
        
        if es_escaneado and _OCR_MODE and OCR_DISPONIBLE:
            print(f"   🔍 OCR activo para PDF escaneado...")
            _stats.pdfs_escaneados += 1
            # Implementación OCR mínima para v6
            doc = fitz.open(file_path)
            pages = []
            for p in doc:
                pix = p.get_pixmap(dpi=300)
                img = Image.open(tempfile.NamedTemporaryFile(suffix=".png"))
                pix.save(img.name)
                pages.append(pytesseract.image_to_string(Image.open(img.name), lang='spa'))
            doc.close()
            _stats.archivos_con_ocr += 1
            return "\n\n".join(pages)
            
        return pymupdf4llm.to_markdown(str(file_path))
    except Exception as e:
        if _logger: _logger.error(f"PDF {file_path.name}: {e}")
        return None

def handle_docx(file_path: Path) -> Optional[str]:
    if not DOCX_DISPONIBLE: return None
    try:
        doc = Document(file_path)
        lines = []
        for p in doc.paragraphs:
            text = p.text.strip()
            if not text: continue
            if p.style.name.startswith('Heading'):
                lines.append(f"{'#' * 2} {text}")
            else:
                lines.append(text)
        for t in doc.tables:
            _stats.tablas_detectadas += 1
            lines.append("")
            for r in t.rows:
                cells = [c.text.strip() for c in r.cells]
                lines.append("| " + " | ".join(cells) + " |")
            lines.append("")
            _stats.tablas_convertidas += 1
        return "\n".join(lines)
    except Exception as e:
        return None

def handle_excel(file_path: Path) -> Optional[str]:
    if not PANDAS_DISPONIBLE: return None
    try:
        xl = pd.ExcelFile(file_path)
        md = []
        for sn in xl.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sn)
            md.append(f"## {sn}\n")
            if not df.empty:
                _stats.tablas_detectadas += 1
                md.append(df.to_markdown(index=False))
                _stats.tablas_convertidas += 1
        return "\n".join(md)
    except Exception as e:
        return None

def handle_pptx(file_path: Path) -> Optional[str]:
    if not PPTX_DISPONIBLE: return None
    try:
        prs = Presentation(file_path)
        md = []
        for i, slide in enumerate(prs.slides, 1):
            md.append(f"## Diapositiva {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text"): md.append(shape.text.strip())
                if shape.has_table:
                    _stats.tablas_detectadas += 1
                    # Conversión básica de tabla pptx
                    for row in shape.table.rows:
                        cells = [c.text_frame.text.strip() for c in row.cells]
                        md.append("| " + " | ".join(cells) + " |")
                    _stats.tablas_convertidas += 1
        return "\n".join(md)
    except Exception as e:
        return None

def handle_png(file_path: Path) -> Optional[str]:
    if not OCR_DISPONIBLE: return f"![{file_path.stem}](./{file_path.name})"
    try:
        _stats.archivos_con_ocr += 1
        return pytesseract.image_to_string(Image.open(file_path), lang='spa')
    except:
        return f"![{file_path.stem}](./{file_path.name})"

def handle_txt(file_path: Path) -> Optional[str]:
    try: return file_path.read_text(encoding='utf-8')
    except: return file_path.read_text(encoding='latin-1', errors='ignore')


# ===========================================================================
# DISPATCHER & CORE
# ===========================================================================

HANDLERS = {
    '.pdf': handle_pdf, '.docx': handle_docx,
    '.xlsx': handle_excel, '.xls': handle_excel, '.xlsm': handle_excel,
    '.pptx': handle_pptx, '.txt': handle_txt, '.png': handle_png,
}

def convertir_archivo(file_path: Path, output_path: Path, verbose: bool = True) -> bool:
    t_start = time.time()
    ext = file_path.suffix.lower()
    
    if ext not in HANDLERS: return False
    
    try:
        if verbose: print(f"\n📖 v6 Procesando: {file_path.name}")
        
        raw = HANDLERS[ext](file_path)
        if raw is None: return False
        
        limpio = limpiar_markdown_completo(raw, verbose=verbose)
        final = get_frontmatter(file_path) + limpio
        
        out_name = normalizar_nombre_archivo(output_path.name)
        out_final = output_path.parent / out_name
        
        out_final.parent.mkdir(parents=True, exist_ok=True)
        out_final.write_text(final, encoding='utf-8')
        
        _stats.incrementar_tipo(ext)
        if verbose: print(f"   ✅ Éxito: {out_final.name} ({time.time()-t_start:.1f}s)")
        return True
    except Exception as e:
        if _logger: _logger.error(f"{file_path.name}: {e}")
        return False

def procesar_carpeta(input_dir: Path, output_dir: Path, workers: int = 1):
    setup_logger(output_dir)
    archivos = []
    for ext in HANDLERS.keys():
        archivos.extend(input_dir.rglob(f"*{ext}"))
        archivos.extend(input_dir.rglob(f"*{ext.upper()}"))
    
    archivos = list(set(archivos))
    if not archivos: return print("❌ No hay archivos.")
    
    print(f"🚀 v6 Iniciando proceso para {len(archivos)} archivos...")
    t_inicio = time.time()
    
    with ThreadPoolExecutor(max_workers=workers) as executor:
        for f in archivos:
            executor.submit(convertir_archivo, f, output_dir / f.relative_to(input_dir).with_suffix('.md'), verbose=(workers==1))
            
    _stats.tiempo_total = time.time() - t_inicio
    _stats.mostrar_resumen()

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Conversor v6 - Tabla Protegida")
    parser.add_argument("--input", default="pdf")
    parser.add_argument("--output", default="docs")
    parser.add_argument("--workers", type=int, default=1)
    parser.add_argument("--ocr", action="store_true")
    parser.add_argument("--debug", action="store_true")
    parser.add_argument("--conservar-tablas", action="store_true")
    
    args = parser.parse_args()
    set_modes(debug=args.debug, ocr=args.ocr, strict=args.conservar_tablas)
    _stats = EstadisticasConversion()
    
    procesar_carpeta(Path(args.input), Path(args.output), args.workers)
