#!/usr/bin/env python3
"""
Multi-Format Converter for Docusaurus v7.2
Converts PDF, DOCX, XLSX, PPTX, TXT, PNG to Markdown
100% FREE - No cloud APIs required

Features:
- 🛡 Table Protection: Isolates tables from cleaning pipeline
- 🧹 Signature Cleaning: Removes digital certificate metadata
- 🖼 Image Handling: Saves to static/img and links correctly
- ⚖ Text Justification: CSS justification for paragraphs
- 🐛 MDX Safe: Auto-escapes < and { for Docusaurus
"""

# ===========================================================================
# IMPORTS
# ===========================================================================

import re
import time
import unicodedata
import logging
import tempfile
import shutil
import os
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Optional, Tuple, List
from datetime import datetime
import warnings

# Optional dependencies
try:
    import fitz
    import pymupdf4llm
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    warnings.warn("pymupdf4llm unavailable")

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    warnings.warn("python-docx unavailable")

try:
    import pandas as pd
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False
    warnings.warn("pandas unavailable")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    warnings.warn("python-pptx unavailable")

try:
    from PIL import Image
    import pytesseract
    HAS_OCR = True
except ImportError:
    HAS_OCR = False
    warnings.warn("OCR unavailable")


# ===========================================================================
# CONSTANTS
# ===========================================================================

# Supported file formats
SUPPORTED_FORMATS = {'.pdf', '.docx', '.xlsx', '.xls', '.xlsm', '.pptx', '.txt', '.png'}
IGNORED_FORMATS = {'.exe', '.url'}

# OCR thresholds
PDF_SCANNED_THRESHOLD = 500  # chars/page
PNG_TEXT_THRESHOLD = 50      # min chars to treat as text

# Text markers
SENTENCE_ENDINGS = ('.', ':', ';', '!', '?', ')', ']', '"', "'")
MD_MARKERS = ('#', '-', '*', '|', '>', '```', '<!--')

# Image folder (relative to project root)
STATIC_IMAGE_PATH = Path('static') / 'img'

# Compiled regex patterns for performance
PATTERN_HYPHEN = re.compile(r'([a-záéíóúüñ])-\n([a-záéíóúüñ])', re.IGNORECASE)
PATTERN_TOC_DOTS = re.compile(r'\s*\.{3,}\s*\d+')
PATTERN_HEADER_1 = re.compile(r'^\s*\d+\s*[-–—]\s*[A-Za-z]{2,}.*\d{4,}\s*$')
PATTERN_HEADER_2 = re.compile(r'^[A-Za-z]{2,}.*\d{4,}\s*[-–—]\s*\d+\s*$')
PATTERN_PAGE_NUM = re.compile(r'^\s*\d+\s*$')
PATTERN_SIGNATURE = re.compile(r'`Firmado por[^`]+`', re.IGNORECASE)
PATTERN_CERT_STARS = re.compile(r'\*{2,}\d+\*{2,}')
PATTERN_SIGN_DATE = re.compile(r'el día \d{2}/\d{2}/\d{4}.*$', re.MULTILINE)
PATTERN_TABLE_BLOCK = re.compile(r'(?:^\|.+\n)+', re.MULTILINE)
PATTERN_STRUCTURAL = re.compile(
    r'^(?:#|\-|\*|\d+\.|\||```|<|id:|title:|sidebar_label:|---)',
    re.IGNORECASE
)
PATTERN_IMG_LINK = re.compile(r'!\[(.*?)\]\((.*?)\)')

# Global state
_DEBUG = False
_OCR = False
_STRICT = False
_stats = None
_logger = None


# ===========================================================================
# STATISTICS
# ===========================================================================

class Stats:
    """Conversion statistics tracker."""
    
    def __init__(self):
        self.files = {}
        self.chars_raw = 0
        self.chars_clean = 0
        self.signatures = 0
        self.ocr_count = 0
        self.images_moved = 0
        self.time = 0.0
    
    def add_file(self, ext: str):
        self.files[ext] = self.files.get(ext, 0) + 1
    
    def show(self):
        if not _DEBUG:
            return
        
        print("\n" + "─" * 60)
        print("📊 CONVERSION STATISTICS (v7.2)")
        print("─" * 60)
        
        if self.files:
            print("📁 Files processed:")
            for ext, count in sorted(self.files.items()):
                print(f"   {ext:8s} : {count:3d} files")
        
        if self.chars_raw > 0:
            reduction = self.chars_raw - self.chars_clean
            pct = (reduction / self.chars_raw * 100) if self.chars_raw else 0
            print(f"\n📝 Characters:")
            print(f"   Raw   : {self.chars_raw:,}")
            print(f"   Clean : {self.chars_clean:,}")
            print(f"   Saved : {reduction:,} ({pct:.1f}%)")
        
        if self.signatures > 0:
            print(f"\n✍️  Signatures cleaned: {self.signatures}")
        
        if self.ocr_count > 0:
            print(f"🔍 OCR processed: {self.ocr_count}")
        
        if self.images_moved > 0:
            print(f"🖼️  Images moved: {self.images_moved}")
        
        if self.time > 0:
            print(f"\n⏱  Total time: {self.time:.2f}s")
        
        print("─" * 60 + "\n")


# ===========================================================================
# UTILITIES
# ===========================================================================

def debug(msg: str):
    """Print debug message if debug mode is active."""
    if _DEBUG:
        print(f"[DEBUG] {msg}")


def setup_logger(output_dir: Path) -> Path:
    """Setup error logger."""
    global _logger
    log_file = output_dir / "errors.log"
    
    _logger = logging.getLogger('converter')
    _logger.setLevel(logging.ERROR)
    
    handler = logging.FileHandler(log_file, encoding='utf-8')
    handler.setFormatter(logging.Formatter(
        '%(asctime)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    ))
    
    _logger.addHandler(handler)
    return log_file


def normalize_filename(name: str) -> str:
    """Convert filename to safe slug."""
    stem = Path(name).stem
    clean = unicodedata.normalize("NFKD", stem).encode("ascii", "ignore").decode("ascii")
    clean = re.sub(r'[^\w\s\-]', '', clean).strip()
    clean = re.sub(r'[\s_]+', '-', clean.lower())
    return (clean or 'document') + '.md'


def create_frontmatter(path: Path) -> str:
    """Generate Docusaurus frontmatter."""
    doc_id = re.sub(r'[^\w\-]', '-', path.stem.lower()).strip('-')
    title = path.stem.replace('-', ' ').replace('_', ' ').title()
    return f"""---
id: {doc_id}
title: "{title}"
sidebar_label: "{title}"
---

"""


# ===========================================================================
# TABLE PROTECTION
# ===========================================================================

def isolate_tables(text: str) -> Tuple[str, List[str]]:
    """Extract tables and replace with placeholders."""
    tables = []
    
    def replacer(match):
        tables.append(match.group(0))
        return f"\n<!-- TABLE_{len(tables)-1} -->\n"
    
    clean_text = PATTERN_TABLE_BLOCK.sub(replacer, text)
    debug(f"Isolated {len(tables)} tables")
    return clean_text, tables


def restore_tables(text: str, tables: List[str]) -> str:
    """Restore tables at placeholder positions."""
    for idx, table in enumerate(tables):
        text = text.replace(f"<!-- TABLE_{idx} -->", f"\n{table}\n")
    return text


# ===========================================================================
# MARKDOWN CLEANING
# ===========================================================================

def clean_signatures(text: str) -> str:
    """Remove digital signature metadata."""
    original = len(text)
    
    text = PATTERN_SIGNATURE.sub('[DIGITAL SIGNATURE]', text)
    text = PATTERN_CERT_STARS.sub('', text)
    text = PATTERN_SIGN_DATE.sub('', text)
    
    if len(text) < original:
        _stats.signatures += 1
    
    return text


def fix_mdx(text: str) -> str:
    """Make content MDX-safe."""
    # Self-closing tags
    text = re.sub(r'<br\s*>', '<br />', text, flags=re.IGNORECASE)
    text = re.sub(r'<hr\s*>', '<hr />', text, flags=re.IGNORECASE)
    
    # Escape < that aren't valid HTML tags
    text = re.sub(
        r'<(?!/?(?:p|br|hr|img|div|span|strong|em|b|i|u|h\d|table|tr|td|th)\b)',
        '&lt;',
        text,
        flags=re.IGNORECASE
    )
    
    # Escape curly braces (JSX conflicts)
    text = text.replace('{', '&#123;').replace('}', '&#125;')
    
    return text


def clean_text(text: str) -> str:
    """Apply cleaning pipeline."""
    _stats.chars_raw += len(text)
    
    # 1. Clean signatures
    text = clean_signatures(text)
    
    # 2. Fix MDX
    text = fix_mdx(text)
    
    # 3. Isolate tables
    text, tables = isolate_tables(text)
    
    # 4. Join hyphenated words
    text = PATTERN_HYPHEN.sub(r'\1\2', text)
    
    # 5. Clean TOC dots
    text = PATTERN_TOC_DOTS.sub('', text)
    
    # 6. Remove headers/footers
    lines = text.split('\n')
    lines = [l for l in lines if not any(p.match(l) for p in [
        PATTERN_HEADER_1, PATTERN_HEADER_2, PATTERN_PAGE_NUM
    ])]
    text = '\n'.join(lines)
    
    # 7. Merge broken paragraphs (unless strict mode)
    if not _STRICT:
        lines = text.split('\n')
        merged = []
        i = 0
        
        while i < len(lines):
            line = lines[i].rstrip()
            
            if (not line or 
                line.startswith(MD_MARKERS) or 
                line.endswith(SENTENCE_ENDINGS)):
                merged.append(line)
                i += 1
                continue
            
            if i + 1 < len(lines):
                next_line = lines[i + 1].lstrip()
                if next_line and not next_line.startswith(MD_MARKERS):
                    merged.append(line + ' ' + next_line)
                    i += 2
                    continue
            
            merged.append(line)
            i += 1
        
        text = '\n'.join(merged)
    
    # 8. Normalize whitespace
    text = re.sub(r' {2,}', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 9. Justify paragraphs
    def justify_callback(match):
        block = match.group(0).strip()
        lines = block.split('\n')
        
        if any(PATTERN_STRUCTURAL.match(l.strip()) for l in lines) or len(block) < 60:
            return match.group(0)
        
        return f'\n<div style={{{{textAlign: "justify"}}}}>\n\n{block}\n\n</div>\n'
    
    text = re.sub(r'((?:(?!\n\n).)+)', justify_callback, text, flags=re.DOTALL)
    
    # 10. Restore tables
    text = restore_tables(text, tables)
    
    # 11. Final cleanup
    text = text.strip()
    _stats.chars_clean += len(text)
    
    return text


# ===========================================================================
# FILE HANDLERS
# ===========================================================================

def handle_pdf(path: Path, img_dir: Path) -> Optional[str]:
    """Convert PDF to Markdown."""
    if not HAS_PDF:
        return None
    
    try:
        # Detect if scanned
        doc = fitz.open(path)
        text = ''.join(p.get_text() for p in doc)
        is_scanned = (len(text) / len(doc)) < PDF_SCANNED_THRESHOLD if len(doc) > 0 else True
        doc.close()
        
        # Apply OCR if needed
        if is_scanned and _OCR and HAS_OCR:
            debug(f"OCR: {path.name}")
            _stats.ocr_count += 1
            
            doc = fitz.open(path)
            pages = []
            
            for page in doc:
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    page.get_pixmap(dpi=300).save(tmp.name)
                    pages.append(pytesseract.image_to_string(Image.open(tmp.name), lang='spa'))
                    os.unlink(tmp.name)
            
            doc.close()
            return '\n\n'.join(pages)
        
        # Normal extraction with image handling
        doc_slug = re.sub(r'[^\w\-]', '-', path.stem.lower())
        doc_img_dir = img_dir / doc_slug
        doc_img_dir.mkdir(parents=True, exist_ok=True)
        
        markdown = pymupdf4llm.to_markdown(
            str(path),
            write_images=True,
            image_path=str(doc_img_dir)
        )
        
        # Fix image paths for Docusaurus
        def fix_img_path(match):
            img_name = Path(match.group(2)).name
            _stats.images_moved += 1
            return f'![{match.group(1)}](/img/{doc_slug}/{img_name})'
        
        return PATTERN_IMG_LINK.sub(fix_img_path, markdown)
    
    except Exception as e:
        if _logger:
            _logger.error(f"PDF {path.name}: {e}")
        return None


def handle_docx(path: Path) -> Optional[str]:
    """Convert DOCX to Markdown."""
    if not HAS_DOCX:
        return None
    
    try:
        doc = Document(path)
        lines = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            if para.style.name.startswith('Heading'):
                lines.append(f"## {text}")
            else:
                lines.append(text)
        
        for table in doc.tables:
            lines.append('')
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append('| ' + ' | '.join(cells) + ' |')
            lines.append('')
        
        return '\n'.join(lines)
    
    except Exception as e:
        if _logger:
            _logger.error(f"DOCX {path.name}: {e}")
        return None


def handle_excel(path: Path) -> Optional[str]:
    """Convert Excel to Markdown."""
    if not HAS_EXCEL:
        return None
    
    try:
        excel_file = pd.ExcelFile(path)
        sections = []
        
        for sheet in excel_file.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet)
            sections.append(f"## {sheet}\n")
            if not df.empty:
                sections.append(df.to_markdown(index=False))
        
        return '\n'.join(sections)
    
    except Exception as e:
        if _logger:
            _logger.error(f"Excel {path.name}: {e}")
        return None


def handle_pptx(path: Path) -> Optional[str]:
    """Convert PowerPoint to Markdown."""
    if not HAS_PPTX:
        return None
    
    try:
        prs = Presentation(path)
        lines = []
        
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n## Slide {i}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    lines.append(shape.text.strip())
                
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text_frame.text.strip() for c in row.cells]
                        lines.append('| ' + ' | '.join(cells) + ' |')
        
        return '\n'.join(lines)
    
    except Exception as e:
        if _logger:
            _logger.error(f"PPTX {path.name}: {e}")
        return None


def handle_txt(path: Path) -> Optional[str]:
    """Read text file."""
    try:
        return path.read_text(encoding='utf-8')
    except UnicodeDecodeError:
        return path.read_text(encoding='latin-1', errors='ignore')


def handle_png(path: Path, img_dir: Path) -> Optional[str]:
    """Handle PNG images with OCR or save to static."""
    # Always save to static/img
    img_name = path.name
    dest = img_dir / img_name
    dest.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(path, dest)
    _stats.images_moved += 1
    
    # Try OCR if available
    if _OCR and HAS_OCR:
        try:
            text = pytesseract.image_to_string(Image.open(path), lang='spa').strip()
            if len(text) > PNG_TEXT_THRESHOLD:
                _stats.ocr_count += 1
                return f"{text}\n\n![Image]({img_name})"
        except:
            pass
    
    return f"![{path.stem}](/img/{img_name})"


# Handler mapping
HANDLERS = {
    '.pdf': handle_pdf,
    '.docx': handle_docx,
    '.xlsx': handle_excel,
    '.xls': handle_excel,
    '.xlsm': handle_excel,
    '.pptx': handle_pptx,
    '.txt': handle_txt,
}


# ===========================================================================
# MAIN PROCESSING
# ===========================================================================

def convert_file(path: Path, output_dir: Path, img_dir: Path) -> bool:
    """Convert single file."""
    ext = path.suffix.lower()
    
    if ext not in HANDLERS and ext != '.png':
        return False
    
    try:
        # Get raw content
        if ext == '.png':
            raw = handle_png(path, img_dir)
        elif ext == '.pdf':
            raw = handle_pdf(path, img_dir)
        else:
            raw = HANDLERS[ext](path)
        
        if not raw:
            return False
        
        # Clean content
        clean = clean_text(raw)
        
        # Add frontmatter
        content = create_frontmatter(path) + clean
        
        # Save
        output_file = output_dir / normalize_filename(path.name)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        output_file.write_text(content, encoding='utf-8')
        
        _stats.add_file(ext)
        debug(f"✓ {path.name}")
        return True
    
    except Exception as e:
        if _logger:
            _logger.error(f"{path.name}: {e}")
        debug(f"✗ {path.name}: {e}")
        return False


def process_all(input_dir: Path, output_dir: Path, workers: int = 1):
    """Process all files in directory."""
    global _stats
    _stats = Stats()
    
    # Setup logging
    setup_logger(output_dir)
    
    # Determine image directory (static/img relative to project root)
    if output_dir.name == 'docs':
        project_root = output_dir.parent
    else:
        project_root = output_dir.resolve()
    
    img_dir = project_root / STATIC_IMAGE_PATH
    img_dir.mkdir(parents=True, exist_ok=True)
    
    # Find files
    files = [
        f for f in input_dir.rglob('*')
        if f.suffix.lower() in SUPPORTED_FORMATS
    ]
    
    if not files:
        print("❌ No files found")
        return
    
    print(f"🚀 v7.2 Processing {len(files)} files...")
    print(f"📂 Output: {output_dir}")
    print(f"🖼️  Images: {img_dir}\n")
    
    t_start = time.time()
    
    # Process
    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = {
            executor.submit(convert_file, f, output_dir, img_dir): f
            for f in files
        }
        
        success = 0
        failed = 0
        
        for future in as_completed(futures):
            if future.result():
                success += 1
            else:
                failed += 1
    
    _stats.time = time.time() - t_start
    
    # Summary
    print(f"\n{'='*60}")
    print(f"✅ Success: {success}")
    print(f"❌ Failed : {failed}")
    print(f"⏱  Time   : {_stats.time:.1f}s")
    print(f"{'='*60}\n")
    
    _stats.show()


# ===========================================================================
# CLI
# ===========================================================================

if __name__ == '__main__':
    import argparse
    
    parser = argparse.ArgumentParser(
        description='Multi-Format Converter for Docusaurus v7.2',
        epilog='Examples:\n'
               '  python converter.py\n'
               '  python converter.py --ocr --workers 4\n'
               '  python converter.py --debug --strict-tables\n',
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    
    parser.add_argument('--input', default='pdf', help='Input folder (default: pdf/)')
    parser.add_argument('--output', default='docs', help='Output folder (default: docs/)')
    parser.add_argument('--workers', type=int, default=1, help='Parallel workers (default: 1)')
    parser.add_argument('--ocr', action='store_true', help='Enable OCR for scanned PDFs')
    parser.add_argument('--debug', action='store_true', help='Debug mode')
    parser.add_argument('--strict-tables', action='store_true', help='Conservative paragraph merging')
    
    args = parser.parse_args()
    
    # Set modes
    _DEBUG = args.debug
    _OCR = args.ocr
    _STRICT = args.strict_tables
    
    # Run
    process_all(Path(args.input), Path(args.output), args.workers)
